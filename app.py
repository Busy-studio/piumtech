import os
import re
import json
import math
import base64
import tempfile
import zipfile
import shutil
import requests
from contextlib import ExitStack
from io import BytesIO
from typing import Any, Dict, List, Tuple

import fitz
import streamlit as st
from PIL import Image, ImageDraw, ImageFont, ImageStat, ImageEnhance, ImageFilter, ImageOps
import matplotlib
import numpy as np
matplotlib.use("Agg")
# Premium infographic style reference asset
PREMIUM_REFERENCE_IMAGE_PATH = "assets/reference/premium_infographic_reference.png"

import matplotlib.pyplot as plt
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

TEXT_MODEL_FIXED = "gpt-4.1-mini"

# Image model routing
# - Asset/detail images for normal SMK sections: lightweight model
# - Final premium reference-based infographic: higher-quality model
ASSET_IMAGE_MODEL_FIXED = os.getenv("ASSET_IMAGE_MODEL", "gpt-image-1-mini")
PREMIUM_IMAGE_MODEL_FIXED = os.getenv("PREMIUM_IMAGE_MODEL", "gpt-image-2")

# Legacy alias retained only for any older helper that still references IMAGE_MODEL_FIXED.
IMAGE_MODEL_FIXED = ASSET_IMAGE_MODEL_FIXED

LANGUAGE_OPTIONS = {
    "한국어": "ko",
    "English": "en",
    "中文": "zh",
    "日本語": "ja",
}

LABELS = {
    "ko": {
        "apps": "적용분야 / 제품",
        "market": "시장현황",
        "overview": "기술개요",
        "diff": "핵심 차별성",
        "drawing": "대표도면",
        "competitiveness": "기술 경쟁력",
        "limitations": "기존기술 한계",
        "advantages": "기술적 우위",
        "ip": "지식재산권 현황",
        "contact": "문의처",
        "invention": "발명의 명칭",
        "app_no": "출원번호\n(등록번호)",
        "app_date": "출원일자\n(등록일자)",
        "prof_suffix": "교수",
        "source": "출처",
        "market_default": "글로벌 시장현황",
    },
    "en": {
        "apps": "Applications / Products",
        "market": "Market Overview",
        "overview": "Technology Overview",
        "diff": "Key Differentiators",
        "drawing": "Representative Drawing",
        "competitiveness": "Technical Competitiveness",
        "limitations": "Current Limitations",
        "advantages": "Technical Advantages",
        "ip": "IP Status",
        "contact": "Contact",
        "invention": "Invention Title",
        "app_no": "Application No.\n(Registration No.)",
        "app_date": "Application Date\n(Registration Date)",
        "prof_suffix": "Professor",
        "source": "Source",
        "market_default": "Global Market Overview",
    },
    "zh": {
        "apps": "应用领域 / 产品",
        "market": "市场现状",
        "overview": "技术概要",
        "diff": "核心差异化",
        "drawing": "代表图",
        "competitiveness": "技术竞争力",
        "limitations": "现有技术局限",
        "advantages": "技术优势",
        "ip": "知识产权现状",
        "contact": "联系方式",
        "invention": "发明名称",
        "app_no": "申请号\n(登记号)",
        "app_date": "申请日期\n(登记日期)",
        "prof_suffix": "教授",
        "source": "来源",
        "market_default": "全球市场现状",
    },
    "ja": {
        "apps": "適用分野 / 製品",
        "market": "市場動向",
        "overview": "技術概要",
        "diff": "主な差別化要素",
        "drawing": "代表図面",
        "competitiveness": "技術競争力",
        "limitations": "既存技術の限界",
        "advantages": "技術的優位性",
        "ip": "知的財産権の状況",
        "contact": "お問い合わせ",
        "invention": "発明の名称",
        "app_no": "出願番号\n(登録番号)",
        "app_date": "出願日\n(登録日)",
        "prof_suffix": "教授",
        "source": "出典",
        "market_default": "グローバル市場動向",
    },
}

def get_lang_code(value: str | None) -> str:
    if not value:
        return "ko"
    return LANGUAGE_OPTIONS.get(str(value), str(value) if str(value) in ("ko", "en", "zh", "ja") else "ko")

def label(lang: str, key: str) -> str:
    return LABELS.get(get_lang_code(lang), LABELS["ko"]).get(key, key)


def ip_table_label(lang: str, key: str, ip: Dict[str, Any]) -> str:
    """등록정보가 없으면 출원번호/출원일자만 표시한다."""
    has_registration = bool((ip or {}).get("registration_number") or (ip or {}).get("registration_date"))
    if has_registration:
        return label(lang, key)
    plain = {
        "ko": {"app_no": "출원번호", "app_date": "출원일자"},
        "en": {"app_no": "Application No.", "app_date": "Application Date"},
        "zh": {"app_no": "申请号", "app_date": "申请日期"},
        "ja": {"app_no": "出願番号", "app_date": "出願日"},
    }
    return plain.get(get_lang_code(lang), plain["ko"]).get(key, label(lang, key))

def has_hangul_text(s: str) -> bool:
    return bool(re.search(r"[가-힣]", str(s or "")))

UNIVERSITY_TRANSLATIONS = {
    "부산대학교": {"en": "Pusan National University", "zh": "釜山国立大学", "ja": "釜山大学"},
    "국립부경대학교": {"en": "Pukyong National University", "zh": "国立釜庆大学", "ja": "国立釜慶大学"},
    "국립한국해양대학교": {"en": "Korea Maritime & Ocean University", "zh": "国立韩国海洋大学", "ja": "国立韓国海洋大学"},
    "동아대학교": {"en": "Dong-A University", "zh": "东亚大学", "ja": "東亜大学"},
    "동의대학교": {"en": "Dong-Eui University", "zh": "东义大学", "ja": "東義大学"},
    "동서대학교": {"en": "Dongseo University", "zh": "东西大学", "ja": "東西大学"},
    "동명대학교": {"en": "Tongmyong University", "zh": "东明大学", "ja": "東明大学"},
    "신라대학교": {"en": "Silla University", "zh": "新罗大学", "ja": "新羅大学"},
    "울산대학교": {"en": "University of Ulsan", "zh": "蔚山大学", "ja": "蔚山大学"},
    "경남대학교": {"en": "Kyungnam University", "zh": "庆南大学", "ja": "慶南大学"},
    "경상국립대학교": {"en": "Gyeongsang National University", "zh": "庆尚国立大学", "ja": "慶尚国立大学"},
    "국립창원대학교": {"en": "Changwon National University", "zh": "国立昌原大学", "ja": "国立昌原大学"},
    "인제대학교": {"en": "Inje University", "zh": "仁济大学", "ja": "仁済大学"},
}

LANGUAGE_NAMES = {"ko": "Korean", "en": "English", "zh": "Simplified Chinese", "ja": "Japanese"}

def get_display_university(university: str, lang: str) -> str:
    lang = get_lang_code(lang)
    if lang == "ko":
        return str(university or "")
    return UNIVERSITY_TRANSLATIONS.get(str(university or ""), {}).get(lang, str(university or ""))

def convert_korean_date(text: str, lang: str) -> str:
    """YYYY년MM월DD일 형식의 날짜를 선택 언어 형식으로 변환."""
    lang = get_lang_code(lang)
    s = str(text or "")
    if not s:
        return ""
    months_en = ["", "Jan.", "Feb.", "Mar.", "Apr.", "May", "Jun.", "Jul.", "Aug.", "Sep.", "Oct.", "Nov.", "Dec."]
    def repl(m):
        y, mo, d = int(m.group(1)), int(m.group(2)), int(m.group(3))
        if lang == "en":
            return f"{months_en[mo]} {d}, {y}"
        if lang == "zh":
            return f"{y}年{mo}月{d}日"
        if lang == "ja":
            return f"{y}年{mo}月{d}日"
        return f"{y}년{mo:02d}월{d:02d}일"
    return re.sub(r"(\d{4})\s*년\s*(\d{1,2})\s*월\s*(\d{1,2})\s*일", repl, s)

def translate_terms_with_gpt(terms: Dict[str, str], lang: str, preserve_values: List[str] | None = None) -> Dict[str, str]:
    """짧은 UI/메타데이터 문구를 선택 언어로 번역. preserve_values는 절대 번역하지 않음."""
    lang = get_lang_code(lang)
    if lang == "ko":
        return {k: str(v or "") for k, v in terms.items()}
    clean_terms = {k: str(v or "") for k, v in terms.items()}
    if not any(v.strip() for v in clean_terms.values()):
        return clean_terms
    preserve_values = [str(v) for v in (preserve_values or []) if str(v).strip()]
    try:
        client = get_client()
        prompt = f"""
Translate the following SMK metadata fields into {LANGUAGE_NAMES.get(lang, lang)}.
Rules:
- Return JSON only with the same keys.
- Do not translate or alter these exact names/strings if they appear: {preserve_values}
- Translate university, department, organization, position, invention title, applicant and short labels naturally.
- Keep patent numbers, phone numbers and emails unchanged.
- Do not add explanations.

Input JSON:
{json.dumps(clean_terms, ensure_ascii=False)}
"""
        res = client.responses.create(model=TEXT_MODEL_FIXED, input=prompt, temperature=0.1)
        parsed = safe_json_parse(res.output_text)
        return {k: str(parsed.get(k, clean_terms.get(k, "")) or "") for k in clean_terms.keys()}
    except Exception:
        return clean_terms

def localize_smk_data(data: Dict[str, Any], lang: str, university: str, department: str, professor: str) -> Dict[str, Any]:
    """출력용 데이터 중 교수명만 보존하고 나머지 메타 필드를 선택 언어로 현지화."""
    lang = get_lang_code(lang)
    data = dict(data or {})
    data["language"] = lang
    ip = normalize_ip(data.get("ip", {}))
    data["university_display"] = get_display_university(university or data.get("university", ""), lang)

    terms = {
        "department_display": department or data.get("department", ""),
        "original_title": data.get("original_title", ""),
        "ip_title": ip.get("title", ""),
        "ip_applicant": ip.get("applicant", ""),
    }
    tr = translate_terms_with_gpt(terms, lang, preserve_values=[professor])
    data["department_display"] = tr.get("department_display", terms["department_display"])
    if tr.get("original_title"):
        data["original_title"] = tr.get("original_title")
    ip["title"] = tr.get("ip_title", ip.get("title", ""))
    ip["applicant"] = tr.get("ip_applicant", ip.get("applicant", ""))
    ip["application_date"] = convert_korean_date(ip.get("application_date", ""), lang)
    ip["registration_date"] = convert_korean_date(ip.get("registration_date", ""), lang)
    data["ip"] = ip
    return data

def build_contact_text(org: str, name: str, position: str, phone: str, email: str, lang: str) -> str:
    """담당자 성명은 그대로 두고 기관/직책만 선택 언어로 현지화."""
    lang = get_lang_code(lang)
    if lang == "ko":
        org_t, pos_t = str(org or ""), str(position or "")
    else:
        tr = translate_terms_with_gpt({"org": org, "position": position}, lang, preserve_values=[name])
        org_t, pos_t = tr.get("org", str(org or "")), tr.get("position", str(position or ""))
    left = " ".join([v for v in [org_t, str(name or ""), pos_t] if str(v).strip()])
    parts = [left]
    if str(phone or "").strip():
        parts.append(str(phone).strip())
    if str(email or "").strip():
        parts.append(str(email).strip())
    return "  |  ".join(parts)

UNIVERSITIES = [
    "부산대학교", "국립부경대학교", "국립한국해양대학교", "동아대학교", "동의대학교", "동서대학교",
    "동명대학교", "신라대학교", "울산대학교", "경남대학교", "경상국립대학교", "국립창원대학교", "인제대학교", "수기입력"
]

st.set_page_config(page_title="PIUM SMK 생성기", page_icon="📄", layout="wide")

# -----------------------------------------------------
# Client / Font
# -----------------------------------------------------
@st.cache_resource
def get_client() -> OpenAI:
    api_key = st.secrets.get("OPENAI_API_KEY", "") or os.getenv("OPENAI_API_KEY", "")
    api_key = str(api_key).strip().replace('"', '').replace("'", "")
    if not api_key.startswith("sk-"):
        st.error("OPENAI_API_KEY가 설정되지 않았습니다. Streamlit Secrets에 API 키를 입력하세요.")
        st.stop()
    return OpenAI(api_key=api_key)

@st.cache_resource
def load_font(size: int, bold: bool = False):
    # CJK 전체(한국어/중국어/일본어)를 지원하는 폰트를 최우선 사용해 □ 깨짐을 방지한다.
    paths = [
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc" if bold else "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSerifCJK-Bold.ttc" if bold else "/usr/share/fonts/opentype/noto/NotoSerifCJK-Regular.ttc",
        "/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf" if bold else "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
        "/usr/share/fonts/truetype/nanum/NanumBarunGothicBold.ttf" if bold else "/usr/share/fonts/truetype/nanum/NanumBarunGothic.ttf",
        "/usr/share/fonts/truetype/nanum/NanumSquareB.ttf" if bold else "/usr/share/fonts/truetype/nanum/NanumSquareR.ttf",
        "/Library/Fonts/AppleGothic.ttf",
        "C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/YuGothB.ttc" if bold else "C:/Windows/Fonts/YuGothR.ttc",
        "C:/Windows/Fonts/malgunbd.ttf" if bold else "C:/Windows/Fonts/malgun.ttf",
    ]
    for p in paths:
        if p and os.path.exists(p):
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


# -----------------------------------------------------
# Logo assets
# -----------------------------------------------------
LOGO_ZIP_PATH = "logo.zip"
LOGO_EXTRACT_DIR = os.path.join(tempfile.gettempdir(), "pium_logo_assets")

UNIV_ALIAS = {
    "경상대학교": "경상국립대학교",
    "창원대학교": "국립창원대학교",
}

def decode_zip_stem(name: str) -> str:
    """logo.zip 내부의 #Uxxxx 형태 파일명을 실제 한글명으로 복원."""
    stem = os.path.splitext(os.path.basename(name))[0]
    def repl(m):
        return chr(int(m.group(1), 16))
    return re.sub(r"#U([0-9A-Fa-f]{4})", repl, stem)

@st.cache_resource
def prepare_logo_assets() -> Dict[str, str]:
    logo_map: Dict[str, str] = {}
    zip_path_candidates = [
        LOGO_ZIP_PATH,
        os.path.join(os.getcwd(), LOGO_ZIP_PATH),
        os.path.join(os.path.dirname(__file__), LOGO_ZIP_PATH) if "__file__" in globals() else LOGO_ZIP_PATH,
    ]
    zip_path = next((z for z in zip_path_candidates if os.path.exists(z)), None)
    if not zip_path:
        return logo_map

    os.makedirs(LOGO_EXTRACT_DIR, exist_ok=True)
    with zipfile.ZipFile(zip_path, "r") as zf:
        for info in zf.infolist():
            if info.is_dir():
                continue
            if not info.filename.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
                continue
            decoded_name = decode_zip_stem(info.filename)
            ext = os.path.splitext(info.filename)[1].lower() or ".png"
            out_path = os.path.join(LOGO_EXTRACT_DIR, decoded_name + ext)
            with zf.open(info) as src, open(out_path, "wb") as dst:
                shutil.copyfileobj(src, dst)
            logo_map[decoded_name] = out_path
    return logo_map

def get_logo_image(name: str) -> Image.Image | None:
    if not name:
        return None
    logos = prepare_logo_assets()
    key = UNIV_ALIAS.get(name, name)
    path = logos.get(key) or logos.get(name)
    if not path or not os.path.exists(path):
        return None
    try:
        return Image.open(path).convert("RGBA")
    except Exception:
        return None

def get_pium_logo_image() -> Image.Image | None:
    return get_logo_image("PIUM")

def get_piumlink_logo_image() -> Image.Image | None:
    """logo.zip 또는 프로젝트 루트에서 PIUMLINK 로고를 불러온다.
    파일명이 PIUMLINK.png, PIUM_LINK.png, 피움링크.png 등으로 들어와도 최대한 탐색한다.
    """
    # 1) logo.zip에서 정확명/느슨한 이름으로 탐색
    logos = prepare_logo_assets()
    for key, path in logos.items():
        norm = re.sub(r"[^0-9a-zA-Z가-힣]", "", str(key)).lower()
        if norm in ("piumlink", "pium링크", "피움링크") or ("pium" in norm and "link" in norm):
            try:
                return Image.open(path).convert("RGBA")
            except Exception:
                pass

    # 2) 프로젝트 루트/앱 폴더 직접 파일 탐색
    roots = [os.getcwd()]
    if "__file__" in globals():
        roots.append(os.path.dirname(__file__))
    candidates = ["PIUMLINK.png", "PIUM_LINK.png", "piumlink.png", "PiumLink.png", "피움링크.png"]
    for root in roots:
        for fname in candidates:
            path = os.path.join(root, fname)
            if os.path.exists(path):
                try:
                    return Image.open(path).convert("RGBA")
                except Exception:
                    pass
    return None

def fit_logo_on_blue(src: Image.Image, size: Tuple[int, int], bg=(0,55,135), padding=16) -> Image.Image:
    """투명 로고를 파란 박스 안에 비율 유지로 삽입."""
    canvas = Image.new("RGBA", size, bg + (255,))
    if src is None:
        return canvas.convert("RGB")
    im = src.copy().convert("RGBA")
    max_w = max(1, size[0] - padding*2)
    max_h = max(1, size[1] - padding*2)
    im.thumbnail((max_w, max_h), Image.LANCZOS)
    x = (size[0] - im.width)//2
    y = (size[1] - im.height)//2
    canvas.alpha_composite(im, (x,y))
    return canvas.convert("RGB")

def make_logo_box_image(logo_img: Image.Image | None, label: str, size=(155,155), bg=(0,55,135)) -> Image.Image:
    box = fit_logo_on_blue(logo_img, size, bg=bg, padding=18)
    if logo_img is None and label:
        d = ImageDraw.Draw(box)
        f = load_font(34, True)
        bbox = d.textbbox((0,0), label, font=f)
        d.text(((size[0]-(bbox[2]-bbox[0]))//2, (size[1]-(bbox[3]-bbox[1]))//2), label, font=f, fill="white")
    return box


# -----------------------------------------------------
# Theme / logo color utilities
# -----------------------------------------------------
def clamp(v: int) -> int:
    return max(0, min(255, int(v)))

def mix(c1: Tuple[int, int, int], c2: Tuple[int, int, int], t: float) -> Tuple[int, int, int]:
    return tuple(clamp(c1[i] * (1 - t) + c2[i] * t) for i in range(3))

def luminance(c: Tuple[int, int, int]) -> float:
    return 0.2126 * c[0] + 0.7152 * c[1] + 0.0722 * c[2]

def ensure_dark(c: Tuple[int, int, int]) -> Tuple[int, int, int]:
    # 본문/제목에 쓸 수 있도록 너무 밝으면 어둡게 보정
    if luminance(c) > 120:
        return mix(c, (0, 35, 85), 0.45)
    return c

def extract_logo_theme(logo_img: Image.Image | None) -> Dict[str, Tuple[int, int, int]]:
    """선택 대학 로고의 주요 색상을 추출해 전체 SMK 색상 테마로 사용."""
    default_primary = (0, 55, 135)
    if logo_img is None:
        primary = default_primary
    else:
        img = logo_img.copy().convert("RGBA")
        img.thumbnail((180, 180), Image.LANCZOS)
        colors = []
        for r, g, b, a in img.getdata():
            if a < 80:
                continue
            # 흰색/회색/검정에 가까운 영역 제외
            mx, mn = max(r, g, b), min(r, g, b)
            sat = mx - mn
            bright = (r + g + b) / 3
            if bright > 232 or bright < 25 or sat < 28:
                continue
            colors.append((r, g, b))
        if not colors:
            primary = default_primary
        else:
            # 가장 많이 등장하는 계열을 안정적으로 잡기 위해 양자화 후 최빈값 사용
            buckets = {}
            for r, g, b in colors:
                key = (round(r / 24) * 24, round(g / 24) * 24, round(b / 24) * 24)
                buckets[key] = buckets.get(key, 0) + 1
            primary = max(buckets.items(), key=lambda x: x[1])[0]
            primary = tuple(clamp(v) for v in primary)
            primary = ensure_dark(primary)

    secondary = mix(primary, (30, 115, 210), 0.25)
    pale = mix(primary, (255, 255, 255), 0.88)
    pale2 = mix(primary, (255, 255, 255), 0.94)
    line = mix(primary, (220, 230, 242), 0.82)
    accent = mix(primary, (0, 175, 190), 0.35)
    return {
        "primary": primary,
        "secondary": secondary,
        "accent": accent,
        "pale": pale,
        "pale2": pale2,
        "line": line,
        "black": (28, 34, 43),
        "gray": (92, 99, 110),
    }

def make_university_logo_box(logo_img: Image.Image | None, label: str, size=(155,155), bg=(235,243,250), primary=(0,55,135)) -> Image.Image:
    """좌측 상단 대학 로고 박스: 선택 대학 로고 계열의 옅은 배경에 로고 삽입."""
    box = Image.new("RGBA", size, bg + (255,))
    if logo_img is not None:
        im = logo_img.copy().convert("RGBA")
        im.thumbnail((size[0]-26, size[1]-26), Image.LANCZOS)
        x = (size[0]-im.width)//2
        y = (size[1]-im.height)//2
        box.alpha_composite(im, (x, y))
    elif label:
        d = ImageDraw.Draw(box)
        f = load_font(24, True)
        lines = label.replace("국립", "국립\n") if len(label) > 5 else label
        bbox = d.multiline_textbbox((0, 0), lines, font=f, spacing=4)
        tx = (size[0] - (bbox[2]-bbox[0])) // 2
        ty = (size[1] - (bbox[3]-bbox[1])) // 2
        d.multiline_text((tx, ty), lines, font=f, fill=primary, align="center", spacing=4)
    return box.convert("RGB")

def make_transparent_logo_canvas(logo_img: Image.Image | None, size=(155,155), padding=8) -> Image.Image:
    """우측 상단 PIUM 로고: 박스 없이 흰/연한 헤더 위에 로고만 배치."""
    canvas = Image.new("RGBA", size, (255,255,255,0))
    if logo_img is None:
        d = ImageDraw.Draw(canvas)
        f = load_font(30, True)
        d.text((20, 55), "PIUM", font=f, fill=(0,55,135,255))
    else:
        im = logo_img.copy().convert("RGBA")
        im.thumbnail((size[0]-padding*2, size[1]-padding*2), Image.LANCZOS)
        x = (size[0]-im.width)//2
        y = (size[1]-im.height)//2
        canvas.alpha_composite(im, (x, y))
    bg = Image.new("RGBA", size, (255,255,255,0))
    bg.alpha_composite(canvas)
    return bg

# -----------------------------------------------------
# PDF extraction
# -----------------------------------------------------
def save_uploaded_file(uploaded_file) -> str:
    suffix = os.path.splitext(uploaded_file.name)[1] or ".pdf"
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.getbuffer())
        return tmp.name

APPLICATION_NO_RE = re.compile(r"(?<!\d)(\d{2}-\d{4}-\d{7})(?!\d)")
DATE_RE = re.compile(r"(?<!\d)((?:19|20)\d{2})\s*(?:[.\-/년])\s*(0?[1-9]|1[0-2])\s*(?:[.\-/월])\s*(0?[1-9]|[12]\d|3[01])\s*(?:일)?")


def _has_searchable_pdf_text(text: str) -> bool:
    """텍스트 레이어가 실제로 존재하는지 확인한다. 공백/기호만 있으면 스캔 PDF로 본다."""
    return bool(re.search(r"[0-9A-Za-z가-힣]", text or ""))


def _ocr_page_with_openai(page: fitz.Page, page_no: int) -> str:
    """텍스트 레이어가 전혀 없는 스캔 PDF에서만 사용하는 OCR fallback."""
    try:
        pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
        b64 = base64.b64encode(pix.tobytes("png")).decode("ascii")
        client = get_client()
        res = client.responses.create(
            model=TEXT_MODEL_FIXED,
            input=[{
                "role": "user",
                "content": [
                    {"type": "input_text", "text": (
                        f"이 이미지는 한국 특허 관련 PDF의 {page_no}페이지다. "
                        "OCR만 수행하라. 보이는 텍스트를 원문 순서대로 최대한 정확히 전사하고, "
                        "출원번호·출원일자·대표도-도N·도면번호의 숫자와 하이픈을 절대 추정하거나 바꾸지 마라. "
                        "설명이나 요약 없이 OCR 텍스트만 출력하라."
                    )},
                    {"type": "input_image", "image_url": f"data:image/png;base64,{b64}"},
                ],
            }],
            temperature=0,
        )
        return (res.output_text or "").strip()
    except Exception:
        return ""


def extract_patent_pages_text(pdf_path: str) -> tuple[list[str], bool]:
    """페이지별 텍스트와 OCR 사용 여부를 반환한다.

    기본은 PDF 텍스트 레이어를 그대로 사용한다. 문서 전체에서 검색 가능한 텍스트가 하나도
    없을 때만 모든 페이지를 OCR한다.
    """
    doc = fitz.open(pdf_path)
    native = [page.get_text("text") or "" for page in doc]
    native_all = "\n".join(native)
    if _has_searchable_pdf_text(native_all):
        doc.close()
        return native, False

    ocr_pages = []
    for i, page in enumerate(doc):
        ocr_pages.append(_ocr_page_with_openai(page, i + 1))
    doc.close()
    return ocr_pages, True


def extract_patent_text(pdf_path: str) -> str:
    pages, _ = extract_patent_pages_text(pdf_path)
    return "\n".join(pages)[:60000]


def _normalize_date_match(m: re.Match | None) -> str:
    if not m:
        return ""
    y, mo, d = m.group(1), int(m.group(2)), int(m.group(3))
    return f"{y}.{mo:02d}.{d:02d}"


def _is_publication_gazette(first_page_text: str) -> bool:
    t = re.sub(r"\s+", "", first_page_text or "")
    return "공개특허공보" in t


def _extract_labeled_application_no(text: str) -> str:
    """출원번호 의미가 명확한 문맥에 붙은 00-0000-0000000 형식만 인정한다."""
    if not text:
        return ""
    for m in APPLICATION_NO_RE.finditer(text):
        left = text[max(0, m.start() - 100):m.start()]
        right = text[m.end():min(len(text), m.end() + 50)]
        context = left + " " + right
        # 단순 '출원' 단어보다 출원번호/특허출원/application number를 우선하고,
        # 일반 '출원'도 번호와 가까운 문맥에서만 허용한다.
        if re.search(r"출\s*원\s*번\s*호|특\s*허\s*출\s*원|application\s*(?:no\.?|number)", context, re.I):
            return m.group(1)
        near = text[max(0, m.start() - 30):m.start()]
        if re.search(r"출\s*원", near):
            return m.group(1)
    return ""


def _extract_date_near_application_no(text: str, app_no: str) -> str:
    if not text or not app_no:
        return ""
    pos = text.find(app_no)
    if pos < 0:
        return ""
    window = text[max(0, pos - 220):min(len(text), pos + len(app_no) + 320)]
    # 출원일자/출원일 표기가 있으면 그 뒤 날짜를 최우선으로 사용한다.
    dm = re.search(r"출\s*원\s*(?:일\s*자|일)\s*[:：]?\s*" + DATE_RE.pattern, window)
    if dm:
        # DATE_RE가 안쪽 캡처를 가지므로 마지막 3개 날짜 그룹을 사용
        groups = dm.groups()
        y, mo, d = groups[-3], int(groups[-2]), int(groups[-1])
        return f"{y}.{mo:02d}.{d:02d}"
    dm2 = DATE_RE.search(window)
    return _normalize_date_match(dm2)


def extract_authoritative_application_metadata(page_texts: list[str]) -> dict:
    """사용자가 정한 출원 메타 권위 규칙을 코드로 적용한다.

    1) 공개공보: 1페이지 공식 메타에서만 추출.
    2) 비공개공보: 전체 PDF에서 출원번호 문맥과 결합된 형식 일치 번호만 사용.
    3) 못 찾으면 빈 값으로 남겨 UI 수기입력 fallback으로 넘긴다.
    """
    first = page_texts[0] if page_texts else ""
    all_text = "\n".join(page_texts or [])
    out = {
        "is_publication_gazette": _is_publication_gazette(first),
        "is_registration_gazette": "등록특허공보" in re.sub(r"\s+", "", first or ""),
        "application_number": "",
        "application_date": "",
    }

    if out["is_publication_gazette"]:
        # 공개공보의 (21) 출원번호 / (22) 출원일자 공식 메타 영역만 사용한다.
        no_m = re.search(
            r"(?:\(\s*21\s*\)|\[\s*21\s*\])?\s*출\s*원\s*번\s*호\s*[:：]?\s*(\d{2}-\d{4}-\d{7})",
            first,
            re.I,
        )
        date_m = re.search(
            r"(?:\(\s*22\s*\)|\[\s*22\s*\])?\s*출\s*원\s*(?:일\s*자|일)\s*[:：]?\s*" + DATE_RE.pattern,
            first,
            re.I,
        )
        if no_m:
            out["application_number"] = no_m.group(1)
        if date_m:
            groups = date_m.groups()
            y, mo, d = groups[-3], int(groups[-2]), int(groups[-1])
            out["application_date"] = f"{y}.{mo:02d}.{d:02d}"
        return out

    app_no = _extract_labeled_application_no(all_text)
    if app_no:
        out["application_number"] = app_no
        out["application_date"] = _extract_date_near_application_no(all_text, app_no)
    return out


def _render_page_image(doc: fitz.Document, page_index: int, zoom: float = 3.2) -> Image.Image:
    page = doc[page_index]
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    return Image.open(BytesIO(pix.tobytes("png"))).convert("RGB")


def _smart_trim_visual(img: Image.Image, threshold: int = 246, padding: int = 12) -> Image.Image:
    """흰 여백을 제거하되, 너무 공격적으로 자르지 않도록 패딩을 남긴다."""
    import numpy as np
    im = img.convert("RGB")
    arr = np.asarray(im)
    # 흰 배경이 아닌 픽셀. 특허 도면의 옅은 회색선도 살리기 위해 threshold를 약간 높게 둠.
    mask = np.any(arr < threshold, axis=2)
    ys, xs = np.where(mask)
    if len(xs) == 0 or len(ys) == 0:
        return im
    x1, x2 = max(xs.min() - padding, 0), min(xs.max() + padding + 1, im.width)
    y1, y2 = max(ys.min() - padding, 0), min(ys.max() + padding + 1, im.height)
    return im.crop((x1, y1, x2, y2))


def _binary_connected_bboxes(mask):
    """작은 ROI용 연결요소 bbox 추출. mask는 bool 2D."""
    import numpy as np
    h, w = mask.shape
    visited = np.zeros_like(mask, dtype=bool)
    bboxes = []
    for yy in range(h):
        xs = np.where(mask[yy] & (~visited[yy]))[0]
        for sx in xs:
            if visited[yy, sx] or not mask[yy, sx]:
                continue
            stack = [(sx, yy)]
            visited[yy, sx] = True
            minx = maxx = sx
            miny = maxy = yy
            count = 0
            while stack:
                x, y = stack.pop()
                count += 1
                if x < minx: minx = x
                if x > maxx: maxx = x
                if y < miny: miny = y
                if y > maxy: maxy = y
                for nx in (x-1, x, x+1):
                    for ny in (y-1, y, y+1):
                        if nx == x and ny == y:
                            continue
                        if 0 <= nx < w and 0 <= ny < h and (not visited[ny, nx]) and mask[ny, nx]:
                            visited[ny, nx] = True
                            stack.append((nx, ny))
            bboxes.append((minx, miny, maxx+1, maxy+1, count))
    return bboxes


def _crop_qr_by_fixed_area(page_img: Image.Image) -> Image.Image | None:
    """등록공보 1페이지 우측 상단의 QR만 안정적으로 크롭한다.
    QR 주변의 공보 구분선/본문 일부가 같이 들어가지 않도록, 우측 상단 ROI에서
    팽창(dilation)된 QR 클러스터만 찾아 다시 크롭한다.
    """
    import numpy as np
    from PIL import ImageFilter

    w, h = page_img.size
    # QR이 위치하는 우측 상단 영역만 넉넉히 가져옴. 이후 클러스터 분석으로 QR만 분리.
    roi = page_img.crop((int(w * 0.78), int(h * 0.010), int(w * 0.992), int(h * 0.125))).convert("RGB")
    arr = np.asarray(roi)
    raw_mask = np.any(arr < 245, axis=2).astype("uint8") * 255

    # QR의 작은 모듈들을 하나의 덩어리로 연결. 긴 수평선은 ratio 조건에서 제외됨.
    mask_img = Image.fromarray(raw_mask, mode="L").filter(ImageFilter.MaxFilter(13))
    mask = np.asarray(mask_img) > 0
    bboxes = _binary_connected_bboxes(mask)

    candidates = []
    for x1, y1, x2, y2, area in bboxes:
        bw, bh = x2 - x1, y2 - y1
        if bw < 45 or bh < 45:
            continue
        ratio = bw / max(1, bh)
        # QR은 정사각형에 가까움. 공보 가로선처럼 긴 요소는 제외.
        if 0.55 <= ratio <= 1.65:
            # 우측 상단에 가까운 큰 정사각형 덩어리를 우선
            score = area + bw * bh * 0.25 + x1 * 0.4 - y1 * 0.15
            candidates.append((score, x1, y1, x2, y2))

    if not candidates:
        # fallback: 기존보다 좁은 좌표. 그래도 실패하면 None.
        fallback = page_img.crop((int(w * 0.84), int(h * 0.018), int(w * 0.975), int(h * 0.105)))
        qr = _smart_trim_visual(fallback, threshold=245, padding=8)
        if qr.width < 35 or qr.height < 35:
            return None
    else:
        _, x1, y1, x2, y2 = max(candidates, key=lambda t: t[0])
        pad = 6
        x1, y1 = max(0, x1 - pad), max(0, y1 - pad)
        x2, y2 = min(roi.width, x2 + pad), min(roi.height, y2 + pad)
        qr = roi.crop((x1, y1, x2, y2))
        qr = _smart_trim_visual(qr, threshold=245, padding=8)

    # QR 표시용 정사각 캔버스
    if qr.width < 35 or qr.height < 35:
        return None
    side = max(qr.width, qr.height)
    canvas = Image.new("RGB", (side, side), "white")
    canvas.paste(qr, ((side - qr.width)//2, (side - qr.height)//2))
    return canvas


def _crop_representative_from_first_page(page_img: Image.Image) -> Image.Image:
    """1페이지 하단의 '대표도' 실제 도면만 크롭한다.
    전체 페이지가 대표도면 박스에 들어가는 문제를 막기 위해, 페이지 하단 영역에서
    픽셀 밀도가 높은 시각 요소 클러스터를 찾아 대표도면으로 사용한다.
    """
    import numpy as np
    w, h = page_img.size

    # 대표도는 국내 공보 1페이지의 요약 아래, 대체로 하단 45% 영역에 위치
    # 페이지 번호/footer는 제외하기 위해 92%까지만 사용
    y0, y1 = int(h * 0.52), int(h * 0.925)
    x0, x1 = int(w * 0.04), int(w * 0.96)
    roi = page_img.crop((x0, y0, x1, y1)).convert("RGB")
    arr = np.asarray(roi)

    # 흰색이 아닌 픽셀 카운트. 텍스트보다 도면 영역이 행/열 방향 밀도가 높음.
    mask = np.any(arr < 245, axis=2)
    row_counts = mask.sum(axis=1)
    col_counts = mask.sum(axis=0)

    # 도면 행 클러스터 탐색: 너무 낮은 텍스트/잡음은 제외
    row_thr = max(6, int(roi.width * 0.018))
    active_rows = row_counts > row_thr

    clusters = []
    start = None
    for i, active in enumerate(active_rows):
        if active and start is None:
            start = i
        elif not active and start is not None:
            if i - start > 12:
                clusters.append((start, i, int(row_counts[start:i].sum())))
            start = None
    if start is not None and len(active_rows) - start > 12:
        clusters.append((start, len(active_rows), int(row_counts[start:].sum())))

    if clusters:
        # 픽셀량이 가장 큰 클러스터를 대표도면으로 판단
        cy1, cy2, _ = max(clusters, key=lambda t: t[2])
        # 위아래 여유 추가
        cy1 = max(0, cy1 - 25)
        cy2 = min(roi.height, cy2 + 25)
        sub = roi.crop((0, cy1, roi.width, cy2))
    else:
        sub = roi

    # 선택된 행 영역 내에서 실제 도면의 좌우 범위 산정
    arr2 = np.asarray(sub)
    mask2 = np.any(arr2 < 245, axis=2)
    col_counts2 = mask2.sum(axis=0)
    col_thr = max(4, int(sub.height * 0.018))
    xs = np.where(col_counts2 > col_thr)[0]
    ys = np.where(mask2.sum(axis=1) > max(4, int(sub.width * 0.01)))[0]

    if len(xs) and len(ys):
        pad = 28
        cx1, cx2 = max(0, xs.min() - pad), min(sub.width, xs.max() + pad + 1)
        cy1b, cy2b = max(0, ys.min() - pad), min(sub.height, ys.max() + pad + 1)
        sub = sub.crop((cx1, cy1b, cx2, cy2b))

    return _smart_trim_visual(sub, threshold=248, padding=16)



def _extract_best_embedded_image_from_page(doc: fitz.Document, page_index: int, prefer_square: bool = False, min_area: int = 15000) -> Image.Image | None:
    """PDF 페이지 내부 이미지 중 실제 도면/QR에 해당하는 이미지를 직접 추출한다.
    렌더링 페이지를 크롭하는 방식보다 전체 페이지가 잘못 들어가는 문제가 적다.
    """
    try:
        page = doc[page_index]
        candidates = []
        for info in page.get_images(full=True):
            xref = info[0]
            try:
                rects = page.get_image_rects(xref)
                img_info = doc.extract_image(xref)
                w, h = int(img_info.get("width", 0)), int(img_info.get("height", 0))
                area = w * h
                if area < min_area:
                    continue
                ratio = w / max(1, h)
                if prefer_square and not (0.75 <= ratio <= 1.35):
                    continue
                if (not prefer_square) and (w < 120 or h < 90):
                    continue
                # 공보 대표도/도면은 보통 페이지 중하단 또는 도면 페이지 본문 영역에 있음.
                pos_score = 0
                if rects:
                    r = rects[0]
                    pos_score = float(r.y0) * 0.15 - float(r.x0) * 0.03
                score = area + pos_score
                candidates.append((score, xref, img_info))
            except Exception:
                continue
        if not candidates:
            return None
        _, xref, img_info = max(candidates, key=lambda t: t[0])
        im = Image.open(BytesIO(img_info["image"])).convert("RGB")
        return _smart_trim_visual(im, threshold=248, padding=12)
    except Exception:
        return None


def _parse_representative_figure_number(first_text: str) -> str | None:
    m = re.search(r"대\s*표\s*도\s*[-–—]?\s*도\s*([0-9]+)", first_text)
    return m.group(1) if m else None

def extract_representative_drawing(pdf_path: str) -> Image.Image:
    """대표도면 추출.
    1순위: 1페이지에 포함된 대표도 이미지 객체를 직접 추출
    2순위: '대 표 도 - 도N' 번호를 읽어 도면 페이지에서 해당 도면 이미지 객체를 추출
    3순위: 기존 렌더링 크롭 fallback
    """
    doc = fitz.open(pdf_path)
    if len(doc) == 0:
        doc.close()
        return Image.new("RGB", (800, 600), "white")

    first_text = doc[0].get_text("text")

    # 1페이지에 대표도가 이미지 객체로 들어있는 경우가 가장 많음.
    first_embedded = _extract_best_embedded_image_from_page(doc, 0, prefer_square=False, min_area=30000)
    if first_embedded is not None and first_embedded.width > 160 and first_embedded.height > 120:
        doc.close()
        return first_embedded

    # 대표도 번호를 기준으로 도면 페이지에서 실제 이미지 객체 추출
    rep_no = _parse_representative_figure_number(first_text)
    if rep_no:
        patterns = [f"도면{rep_no}", f"도면 {rep_no}", f"도 {rep_no}"]
        for i in range(1, len(doc)):
            t = doc[i].get_text("text")
            if any(p in t for p in patterns):
                img = _extract_best_embedded_image_from_page(doc, i, prefer_square=False, min_area=25000)
                if img is not None:
                    doc.close()
                    return img

    # fallback: 도면 페이지 중 큰 이미지 객체 우선 추출
    for i in range(1, len(doc)):
        t = doc[i].get_text("text")
        if "도면" in t or "도 " in t:
            img = _extract_best_embedded_image_from_page(doc, i, prefer_square=False, min_area=25000)
            if img is not None:
                doc.close()
                return img

    # 최후 fallback: 1페이지 하단 크롭
    first_img = _render_page_image(doc, 0, zoom=3.2)
    doc.close()
    return _crop_representative_from_first_page(first_img)


def extract_qr_code_from_first_page(pdf_path: str) -> Image.Image | None:
    """등록공보 1페이지 우측 상단 QR 코드만 추출한다.
    우선 PDF의 작은 정사각형 이미지 객체를 직접 찾고, 실패 시 렌더링 ROI 크롭으로 보정한다.
    """
    try:
        doc = fitz.open(pdf_path)
        if len(doc) == 0:
            doc.close()
            return None
        page = doc[0]
        page_rect = page.rect
        candidates = []
        for info in page.get_images(full=True):
            xref = info[0]
            try:
                img_info = doc.extract_image(xref)
                iw, ih = int(img_info.get("width", 0)), int(img_info.get("height", 0))
                if iw < 35 or ih < 35:
                    continue
                ratio = iw / max(1, ih)
                if not (0.75 <= ratio <= 1.35):
                    continue
                rects = page.get_image_rects(xref)
                if not rects:
                    continue
                r = rects[0]
                # 1페이지 우측 상단 영역의 정사각형 이미지 = QR 가능성 높음
                if r.x0 < page_rect.width * 0.70 or r.y0 > page_rect.height * 0.18:
                    continue
                score = (page_rect.width - r.x0) + (page_rect.height * 0.18 - r.y0) + iw * ih * 0.01
                candidates.append((score, img_info))
            except Exception:
                continue
        if candidates:
            _, img_info = max(candidates, key=lambda t: t[0])
            qr = Image.open(BytesIO(img_info["image"])).convert("RGB")
            qr = _smart_trim_visual(qr, threshold=245, padding=4)
            side = max(qr.width, qr.height)
            canvas = Image.new("RGB", (side, side), "white")
            canvas.paste(qr, ((side-qr.width)//2, (side-qr.height)//2))
            doc.close()
            return canvas

        img = _render_page_image(doc, 0, zoom=3.2)
        doc.close()
        return _crop_qr_by_fixed_area(img)
    except Exception:
        return None

# -----------------------------------------------------
# GPT
# -----------------------------------------------------
def safe_json_parse(text: str) -> Dict[str, Any]:
    text = text.strip()
    text = re.sub(r"^```json", "", text, flags=re.MULTILINE).strip()
    text = re.sub(r"^```", "", text, flags=re.MULTILINE).strip()
    text = re.sub(r"```$", "", text, flags=re.MULTILINE).strip()
    m = re.search(r"\{.*\}", text, re.S)
    if m:
        text = m.group(0)
    return json.loads(text)

def analyze_patent_with_gpt(patent_text: str, university: str, department: str, professor: str, language: str = "ko") -> Dict[str, Any]:
    client = get_client()
    lang = get_lang_code(language)
    lang_instructions = {
        "ko": "모든 SMK 문구를 한국어로 작성하라. 기술개요/차별성/경쟁력은 개조식 문체로 작성하라.",
        "en": "Write all SMK copy in English. Use concise business-style bullet copy.",
        "zh": "请将所有SMK文案用中文撰写。技术概要、差异化和竞争力内容使用简洁的商务要点式表达。",
        "ja": "SMKのすべての文言を日本語で作成すること。技術概要・差別化要素・競争力は簡潔なビジネス向け箇条書きで記述すること。",
    }
    lang_instruction = lang_instructions.get(lang, lang_instructions["ko"]) + " Keep professor names and contact person names exactly as entered by the user."
    prompt = f"""
너는 대학 기술마케팅자료(SMK) 작성 전문가다.
아래 특허 명세서를 바탕으로 카드형 1페이지 기술소개자료에 들어갈 내용을 생성하라.

언어 설정:
- {lang_instruction}
- 단, 교수명은 사용자가 입력한 값 그대로 유지한다: {professor}

작성 기준:
- 해당 분야 4년제 대학 졸업자가 이해할 수 있는 수준
- 기술이전/사업화 검토자가 빠르게 이해할 수 있는 개조식 문체
- 과장 금지, 특허 명세서 근거 기반
- 기술명은 마케팅용 제목으로 자연스럽게 정리하되 원 발명의 핵심을 유지
- 적용분야 제품은 3개, 각 명칭은 짧게
- 기술개요는 3개
- 기술 차별성은 3개
- 기존기술 한계는 2개
- 기술적 우위는 2개
- 지식재산권의 출원번호와 출원일자는 별도 코드가 권위값을 결정하므로 절대 추정하지 말고, 명확히 보이는 경우에만 적는다
- 등록번호/등록일자도 문서에 명시된 경우에만 적고 추정하지 않는다
- 값이 없으면 빈 문자열로 둔다
- 출력은 JSON만

JSON 형식:
{{
  "marketing_title": "",
  "subtitle": "",
  "original_title": "",
  "university": "{university}",
  "department": "{department}",
  "professor": "{professor}",
  "applications": [
    {{"name": "", "description": ""}},
    {{"name": "", "description": ""}},
    {{"name": "", "description": ""}}
  ],
  "overview": ["", "", ""],
  "differentiation": ["", "", ""],
  "limitations": ["", ""],
  "technical_advantages": ["", ""],
  "ip": {{
    "title": "",
    "application_number": "",
    "registration_number": "",
    "application_date": "",
    "registration_date": "",
    "applicant": ""
  }}
}}

대학교: {university}
학과: {department}
교수명: {professor}

특허 명세서:
{patent_text}
"""
    res = client.responses.create(model=TEXT_MODEL_FIXED, input=prompt, temperature=0.2)
    return safe_json_parse(res.output_text)

def rgb_to_hex(color: Tuple[int, int, int]) -> str:
    return "#%02X%02X%02X" % tuple(int(max(0, min(255, c))) for c in color)


def recolor_icon_palette(img: Image.Image, primary: Tuple[int, int, int], accent: Tuple[int, int, int]) -> Image.Image:
    """적용분야 아이콘이 기본 파랑 계열로 생성되더라도 대학 로고 대표색으로 한 번 더 보정."""
    rgba = img.convert("RGBA")
    out = []
    for r, g, b, a in rgba.getdata():
        if a == 0:
            out.append((r, g, b, a))
            continue
        mx, mn = max(r, g, b), min(r, g, b)
        sat = mx - mn
        bright = (r + g + b) / 3.0

        # 흰 배경 / 거의 무채색은 그대로 유지
        if bright > 245 or sat < 22:
            out.append((r, g, b, a))
            continue

        # 원본이 청록 계열이면 accent, 그 외 유채색은 primary 위주로 매핑
        is_cyanish = (g >= r + 8 and b >= r + 8)
        target = accent if is_cyanish else primary

        # 밝기와 채도를 살려 shading 유지
        tone = 0.22 + 0.68 * (bright / 255.0)
        strength = min(0.82, 0.26 + (sat / 255.0) * 0.56)
        tinted = tuple(clamp(255 - (255 - target[i]) * tone) for i in range(3))
        nr, ng, nb = mix((r, g, b), tinted, strength)
        out.append((nr, ng, nb, a))

    rgba.putdata(out)
    return rgba.convert("RGB")


def generate_application_image(title: str, desc: str, university_logo: Image.Image | None = None) -> Image.Image:
    client = get_client()
    theme = extract_logo_theme(university_logo)
    primary = theme["primary"]
    accent = theme["accent"]
    primary_hex = rgb_to_hex(primary)
    accent_hex = rgb_to_hex(accent)
    prompt = f"""
Create ONE icon from a unified premium technology icon set for a Korean university tech-transfer one-page brochure.
Application: {title}
Description: {desc}

Mandatory visual style:
- pure white background only, no black background, no dark vignette, no colored backdrop
- consistent semi-isometric vector-flat illustration style
- use the selected university logo color family as the main accent palette
- primary brand color: {primary_hex}
- secondary brand color: {accent_hex}
- preserve that brand-accented palette instead of default generic blue if the university logo color is different
- allow white and light gray as neutrals only
- same stroke thickness, same lighting direction, same icon scale
- centered object with generous white margin
- professional public-sector technology marketing style
- no text, no letters, no logos, no watermark
"""
    result = client.images.generate(model=ASSET_IMAGE_MODEL_FIXED, prompt=prompt, size="1024x1024")
    img_b64 = result.data[0].b64_json
    img = Image.open(BytesIO(base64.b64decode(img_b64))).convert("RGB")
    img = clean_dark_background(img)
    return recolor_icon_palette(img, primary, accent)


def generate_application_images_set(apps: List[Dict[str, Any]], university_logo: Image.Image | None = None) -> List[Image.Image]:
    """3개 적용분야 아이콘을 한 번에 생성 후 3등분해 그림체를 최대한 통일."""
    client = get_client()
    theme = extract_logo_theme(university_logo)
    primary = theme["primary"]
    accent = theme["accent"]
    primary_hex = rgb_to_hex(primary)
    accent_hex = rgb_to_hex(accent)
    app_text = []
    for idx, app in enumerate(apps[:3], 1):
        app_text.append(f"{idx}. {app.get('name','')} - {app.get('description','')}")
    joined = "\n".join(app_text)
    prompt = f"""
Create a horizontal set of THREE matching application icons for a Korean university technology brief.
The three icons must look like they belong to the exact same icon family.

Applications:
{joined}

Mandatory layout:
- 3 separate icons arranged left, center, right with large white spacing
- pure white background only
- no dividers, no text, no labels, no logos, no watermark

Mandatory unified visual style:
- consistent semi-isometric vector-flat illustration style
- use the selected university logo color family as the main accent palette
- primary brand color: {primary_hex}
- secondary brand color: {accent_hex}
- preserve that brand-accented palette instead of default generic blue if the university logo color is different
- allow white and light gray as neutrals only
- same stroke thickness, same lighting direction, same icon scale
- centered objects, generous white margin
- professional public-sector technology marketing style
"""
    result = client.images.generate(model=ASSET_IMAGE_MODEL_FIXED, prompt=prompt, size="1536x1024")
    img_b64 = result.data[0].b64_json
    sheet = Image.open(BytesIO(base64.b64decode(img_b64))).convert("RGB")
    sheet = clean_dark_background(sheet)
    w, h = sheet.size
    icons = []
    for i in range(3):
        crop = sheet.crop((i*w//3, 0, (i+1)*w//3, h))
        icons.append(recolor_icon_palette(crop, primary, accent))
    return icons

# -----------------------------------------------------
# Image utilities
# -----------------------------------------------------
def clean_dark_background(img: Image.Image) -> Image.Image:
    """검정/짙은 배경이 섞여 나온 경우 흰 배경으로 완화."""
    img = img.convert("RGB")
    w, h = img.size
    border = Image.new("RGB", (w*2 + h*2, 1), "white")
    px = []
    for x in range(w):
        px.append(img.getpixel((x, 0)))
        px.append(img.getpixel((x, h-1)))
    for y in range(h):
        px.append(img.getpixel((0, y)))
        px.append(img.getpixel((w-1, y)))
    avg = sum((r+g+b)/3 for r,g,b in px) / max(len(px), 1)
    if avg > 130:
        return img
    data = []
    for r,g,b in img.getdata():
        brightness = (r+g+b)/3
        # 매우 어두운 배경만 흰색으로 치환, 오브젝트 음영은 최대한 보존
        if brightness < 42 and max(r,g,b) - min(r,g,b) < 35:
            data.append((255,255,255))
        else:
            data.append((r,g,b))
    img.putdata(data)
    return img

def trim_light_margins(src: Image.Image, threshold: int = 248, padding: int = 8) -> Image.Image:
    """흰색 여백을 자동 제거해 아이콘/대표도면이 박스 안에서 작아 보이지 않도록 보정."""
    im = src.copy().convert("RGB")
    w, h = im.size
    pix = im.load()
    xs, ys = [], []
    for y in range(h):
        for x in range(w):
            r, g, b = pix[x, y]
            # 거의 흰 배경과 아주 연한 회색은 여백으로 간주
            if not (r >= threshold and g >= threshold and b >= threshold):
                # 너무 연한 그림자도 일부 포함되도록 보수적으로 판단
                if max(r, g, b) - min(r, g, b) > 8 or (r + g + b) / 3 < threshold - 6:
                    xs.append(x); ys.append(y)
    if not xs or not ys:
        return im
    x1, x2 = max(0, min(xs)-padding), min(w, max(xs)+padding)
    y1, y2 = max(0, min(ys)-padding), min(h, max(ys)+padding)
    if x2 <= x1 or y2 <= y1:
        return im
    return im.crop((x1, y1, x2, y2))

def fit_image(src: Image.Image, size: Tuple[int, int], bg=(255,255,255), trim: bool = True) -> Image.Image:
    im = src.copy().convert("RGB")
    if trim:
        im = trim_light_margins(im)
    im.thumbnail(size, Image.LANCZOS)
    canvas = Image.new("RGB", size, bg)
    x = (size[0] - im.width)//2
    y = (size[1] - im.height)//2
    canvas.paste(im, (x,y))
    return canvas

def _split_token_to_fit(draw: ImageDraw.ImageDraw, token: str, font, max_width: int) -> List[str]:
    chunks, cur = [], ""
    for ch in str(token):
        test = cur + ch
        if draw.textbbox((0, 0), test, font=font)[2] <= max_width or not cur:
            cur = test
        else:
            chunks.append(cur)
            cur = ch
    if cur:
        chunks.append(cur)
    return chunks

def wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int, max_lines: int = None) -> List[str]:
    """한국어/영어 모두 단어 단위 우선 줄바꿈. 긴 단어만 글자 단위로 보정."""
    lines = []
    for raw in str(text or "").split("\n"):
        raw = raw.strip()
        if not raw:
            lines.append("")
            continue
        tokens = re.findall(r"\S+\s*", raw)
        line = ""
        for tok in tokens:
            tok = tok.rstrip()
            candidate = (line + " " + tok).strip() if line else tok
            if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
                line = candidate
                continue
            if line:
                lines.append(line)
                if max_lines and len(lines) >= max_lines:
                    line = ""
                    break
                line = ""
            if draw.textbbox((0, 0), tok, font=font)[2] <= max_width:
                line = tok
            else:
                chunks = _split_token_to_fit(draw, tok, font, max_width)
                for chunk in chunks[:-1]:
                    lines.append(chunk)
                    if max_lines and len(lines) >= max_lines:
                        break
                if max_lines and len(lines) >= max_lines:
                    line = ""
                    break
                line = chunks[-1] if chunks else ""
        if max_lines and len(lines) >= max_lines:
            break
        if line:
            lines.append(line)
    if max_lines and len(lines) > max_lines:
        lines = lines[:max_lines]
    return lines


def draw_wrapped(draw, xy, text, font, fill, max_width, line_gap=7, max_lines=None) -> int:
    x, y = xy
    lines = wrap_text(draw, text, font, max_width, max_lines=max_lines)
    for i, line in enumerate(lines):
        if max_lines and i == max_lines-1 and len(wrap_text(draw, text, font, max_width)) > max_lines:
            line = line[:-1] + "…"
        draw.text((x, y), line, font=font, fill=fill)
        y += getattr(font, "size", 18) + line_gap
    return y


def draw_centered_wrapped(draw, box, text, font, fill, max_lines=None, line_gap=4):
    """주어진 박스 안에서 텍스트를 가로 중앙 정렬로 출력."""
    x1, y1, x2, y2 = box
    max_width = x2 - x1
    lines = wrap_text(draw, str(text or ""), font, max_width, max_lines=max_lines)
    line_h = getattr(font, "size", 18) + line_gap
    total_h = len(lines) * getattr(font, "size", 18) + max(0, len(lines)-1) * line_gap
    y = y1 + max(0, (y2 - y1 - total_h)//2)
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        tw = bbox[2] - bbox[0]
        draw.text((x1 + (max_width - tw)//2, y), line, font=font, fill=fill)
        y += line_h
    return y

def draw_section_title(draw, x, y, title, font, color):
    draw.rounded_rectangle((x, y, x+10, y+34), radius=4, fill=color)
    draw.text((x+22, y+2), title, font=font, fill=color)

def draw_card(draw, xyxy, radius=24, fill=(255,255,255), outline=(220,230,242), width=2):
    draw.rounded_rectangle(xyxy, radius=radius, fill=fill, outline=outline, width=width)

def bullet_list(draw, x, y, items, font, color, max_width, bullet="•", line_gap=7, item_gap=10, max_lines_per_item=3):
    for item in items:
        item = str(item).strip()
        if not item:
            continue
        draw.text((x, y), bullet, font=font, fill=color)
        y = draw_wrapped(draw, (x+24, y), item, font, color, max_width-24, line_gap, max_lines=max_lines_per_item) + item_gap
    return y

# -----------------------------------------------------
# New Box Layout Rendering - boxed v2
# -----------------------------------------------------
def normalize_ip(ip: Dict[str, Any]) -> Dict[str, str]:
    return {
        "title": str(ip.get("title") or ip.get("name") or ""),
        "application_number": str(ip.get("application_number") or ip.get("app_no") or ""),
        "registration_number": str(ip.get("registration_number") or ip.get("reg_no") or ""),
        "application_date": str(ip.get("application_date") or ip.get("app_date") or ""),
        "registration_date": str(ip.get("registration_date") or ip.get("reg_date") or ""),
        "applicant": str(ip.get("applicant") or ""),
    }


def normalize_market_info(info: Dict[str, Any] | None) -> Dict[str, Any]:
    info = info or {}
    return {
        "language": str(info.get("language") or "ko"),
        "market_scope": str(info.get("market_scope") or "global"),
        "market_name": str(info.get("market_name") or info.get("name") or ""),
        "display_title": str(info.get("display_title") or info.get("title") or ""),
        "base_year": str(info.get("base_year") or ""),
        "end_year": str(info.get("end_year") or ""),
        "base_value": str(info.get("base_value") or ""),
        "end_value": str(info.get("end_value") or ""),
        "unit": str(info.get("unit") or ""),
        "cagr": str(info.get("cagr") or info.get("growth_rate") or ""),
        "summary": str(info.get("summary") or ""),
        "source_title": str(info.get("source_title") or ""),
        "source_url": str(info.get("source_url") or ""),
        "source_year": str(info.get("source_year") or ""),
        "graph_image_url": str(info.get("graph_image_url") or ""),
        "graph_image_caption": str(info.get("graph_image_caption") or ""),
        "verified": bool(info.get("verified", False)),
        "error": str(info.get("error") or ""),
    }

def _safe_float(v, default=0.0):
    try:
        s = str(v).replace(",", "").replace("%", "").strip()
        return float(s) if s else default
    except Exception:
        return default

def _safe_int(v, default=0):
    try:
        s = re.sub(r"[^0-9-]", "", str(v))
        return int(s) if s else default
    except Exception:
        return default


def empty_market_info(reason: str = "검증 가능한 글로벌 시장 데이터를 찾지 못했습니다.") -> Dict[str, Any]:
    return normalize_market_info({
        "market_scope": "global",
        "market_name": "글로벌 시장",
        "display_title": "글로벌 시장현황",
        "summary": "검증 가능한 출처 확인 후 입력이 필요합니다.",
        "verified": False,
        "error": reason,
    })

def has_market_graph_image(info: Dict[str, Any]) -> bool:
    m = normalize_market_info(info)
    return bool(m.get("graph_image_url") and m.get("source_title") and m.get("source_url"))

def download_image_from_url(url: str, timeout: int = 12) -> Image.Image | None:
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        r = requests.get(url, timeout=timeout, headers=headers)
        r.raise_for_status()
        ctype = (r.headers.get("content-type") or "").lower()
        if not any(x in ctype for x in ["image/", "application/octet-stream", "binary/"]):
            return None
        return Image.open(BytesIO(r.content)).convert("RGB")
    except Exception:
        return None

def make_market_placeholder(primary=(0,55,135), message1="검증된 시장 그래프 없음", message2="출처 확인 후 시장현황을 업데이트하세요") -> Image.Image:
    img = Image.new("RGB", (920, 310), "white")
    d = ImageDraw.Draw(img)
    d.rectangle((0, 0, 919, 309), outline=(220,230,242), width=2)
    d.text((70, 110), message1, font=load_font(34, True), fill=primary)
    d.text((70, 165), message2, font=load_font(23, False), fill=(92,99,110))
    return img

def get_market_visual(info: Dict[str, Any], primary=(0,55,135), accent=(0,165,185)) -> tuple[Image.Image, str]:
    if is_valid_market_info(info):
        return generate_market_chart(info, primary=primary, accent=accent), "generated_chart"
    if has_market_graph_image(info):
        img = download_image_from_url(normalize_market_info(info).get("graph_image_url", ""))
        if img is not None:
            return img, "source_graph_image"
    return make_market_placeholder(primary=primary), "placeholder"

def is_valid_market_info(info: Dict[str, Any]) -> bool:
    m = normalize_market_info(info)
    has_source = bool(m.get("source_title") and m.get("source_url"))
    has_cagr = bool(str(m.get("cagr") or "").strip())
    has_values = bool(str(m.get("base_value") or "").strip() and str(m.get("end_value") or "").strip())
    has_years = bool(str(m.get("base_year") or "").strip() and str(m.get("end_year") or "").strip())
    return has_source and has_years and (has_cagr or has_values)


def select_market_candidates(data: Dict[str, Any]) -> List[Dict[str, str]]:
    """적용분야를 대표할 수 있는 상위 글로벌 시장 후보 2~3개를 뽑는다."""
    client = get_client()
    apps = data.get("applications", [])[:3]
    app_lines = []
    for i, app in enumerate(apps, 1):
        app_lines.append(f"{i}. {app.get('name','')} - {app.get('description','')}")
    app_text = "\n".join(app_lines) or "적용분야 정보 없음"
    prompt = f"""
너는 대학 기술소개자료(SMK)의 시장현황 구성을 돕는 시장 분류 전문가다.
중요: 기술명 자체를 시장명으로 사용하지 마라.
아래 기술의 적용분야를 보고, 대표성이 있는 상위 글로벌 시장 후보를 2~3개 뽑아라.

규칙:
- 기술명 그대로를 시장명으로 쓰지 않는다.
- 반드시 실제 시장조사 리포트에서 흔히 쓰이는 넓은 시장명으로 정한다.
- 각 후보는 서로 다른 관점의 상위 시장이어야 한다.
- 시장명은 너무 좁은 특허명/솔루션명이 아니라 상위 카테고리여야 한다.
- 예시: Global Remote Sensing Market, Global Geospatial Analytics Market, Global Precision Agriculture Market, Global GIS Market, Global Environmental Monitoring Market
- 출력은 JSON만 반환한다.

JSON 형식:
{{
  "candidates": [
    {{"market_name_ko": "", "market_name_en": "", "display_title": "", "search_keywords": ""}},
    {{"market_name_ko": "", "market_name_en": "", "display_title": "", "search_keywords": ""}},
    {{"market_name_ko": "", "market_name_en": "", "display_title": "", "search_keywords": ""}}
  ]
}}

기술명: {data.get('marketing_title','')}
기술 한줄요약: {data.get('subtitle','')}
적용분야:
{app_text}
"""
    try:
        res = client.responses.create(model=TEXT_MODEL_FIXED, input=prompt, temperature=0.1)
        parsed = safe_json_parse(res.output_text)
        raw_candidates = parsed.get("candidates") or []
        candidates: List[Dict[str, str]] = []
        for item in raw_candidates[:3]:
            if not isinstance(item, dict):
                continue
            ko = str(item.get("market_name_ko") or "").strip()
            en = str(item.get("market_name_en") or "").strip()
            display = str(item.get("display_title") or "").strip()
            keywords = str(item.get("search_keywords") or "").strip()
            if not (ko or en or display or keywords):
                continue
            if not display:
                base = ko or en or "대표 글로벌 시장"
                display = f"글로벌 {base} 시장"
            candidates.append({
                "market_name_ko": ko,
                "market_name_en": en,
                "display_title": display,
                "search_keywords": keywords or (en or ko),
            })
        if candidates:
            return candidates
    except Exception:
        pass

    fallback_name = apps[0].get("name", "대표 적용시장") if apps else "대표 적용시장"
    return [{
        "market_name_ko": fallback_name,
        "market_name_en": fallback_name,
        "display_title": f"글로벌 {fallback_name} 시장",
        "search_keywords": fallback_name,
    }]


def choose_best_market_candidate(data: Dict[str, Any], candidates: List[Dict[str, str]]) -> Dict[str, str]:
    """후보 2~3개 중 공개 글로벌 시장 데이터가 가장 잘 잡히는 후보를 자동 선택한다."""
    if not candidates:
        return {
            "market_name_ko": "대표 적용시장",
            "market_name_en": "Representative Market",
            "display_title": "글로벌 대표 적용시장",
            "search_keywords": "Representative Market",
        }
    if len(candidates) == 1:
        return candidates[0]

    client = get_client()
    apps = data.get("applications", [])[:3]
    app_lines = []
    for i, app in enumerate(apps, 1):
        app_lines.append(f"{i}. {app.get('name','')} - {app.get('description','')}")
    app_text = "\n".join(app_lines) or "적용분야 정보 없음"
    cand_text_lines = []
    for idx, c in enumerate(candidates, 1):
        cand_text_lines.append(
            f"{idx}. 한글명: {c.get('market_name_ko','')} | 영문명: {c.get('market_name_en','')} | 표시제목: {c.get('display_title','')} | 검색키워드: {c.get('search_keywords','')}"
        )
    cand_text = "\n".join(cand_text_lines)
    prompt = f"""
너는 시장 리서치 전문가다.
아래 후보 2~3개 중에서, 공개 웹 검색으로 글로벌 시장 CAGR/시장규모/전망/그래프를 가장 찾기 쉬운 후보 1개를 고르라.

선정 기준:
- Global/Worldwide 시장 데이터가 실제로 많이 공개되어 있을 것
- CAGR 또는 시장규모가 비교적 쉽게 확인될 것
- 시장조사 리포트/공개 페이지/그래프 이미지를 확보하기 쉬울 것
- 기술 적용분야와도 충분히 연관성이 있을 것
- 기술명 자체와 비슷한 좁은 시장 대신 상위 시장을 선호할 것

출력은 JSON만 반환한다.

JSON 형식:
{{
  "selected_index": 1,
  "reason": ""
}}

기술명: {data.get('marketing_title','')}
기술 한줄요약: {data.get('subtitle','')}
적용분야:
{app_text}

시장 후보 목록:
{cand_text}
"""
    try:
        res = client.responses.create(
            model=TEXT_MODEL_FIXED,
            input=prompt,
            tools=[{"type": "web_search_preview"}],
            temperature=0.1,
        )
        parsed = safe_json_parse(res.output_text)
        idx = int(parsed.get("selected_index", 1)) - 1
        if 0 <= idx < len(candidates):
            chosen = dict(candidates[idx])
            chosen["selection_reason"] = str(parsed.get("reason") or "").strip()
            return chosen
    except Exception:
        pass
    return candidates[0]


def select_representative_market(data: Dict[str, Any]) -> Dict[str, str]:
    """후보 2~3개를 뽑고, 그중 검색 결과가 가장 잘 나오는 시장을 최종 선택한다."""
    candidates = select_market_candidates(data)
    return choose_best_market_candidate(data, candidates)


def generate_market_info_with_web(data: Dict[str, Any]) -> Dict[str, Any]:
    client = get_client()
    lang = get_lang_code(data.get("language", "ko"))
    market_lang_rules = {
        "ko": "summary와 display_title은 반드시 한국어로 작성한다. summary는 개조식/음슴체에 어울리는 짧은 표현으로 작성한다.",
        "en": "summary and display_title must be written in English. Use concise business-style wording.",
        "zh": "summary 和 display_title 必须用中文撰写，使用简洁的商务要点式表达。",
        "ja": "summary と display_title は必ず日本語で作成し、簡潔なビジネス向け表現にする。",
    }
    lang_rule = market_lang_rules.get(lang, market_lang_rules["ko"])
    apps = data.get("applications", [])[:3]
    app_lines = []
    for i, app in enumerate(apps, 1):
        app_lines.append(f"{i}. {app.get('name','')} - {app.get('description','')}")
    app_text = "\n".join(app_lines)
    rep_market = select_representative_market(data)

    prompt = f"""
너는 대학 기술소개자료(SMK)에 들어갈 시장현황을 조사하는 리서치 전문가다.
반드시 글로벌(Global/Worldwide) 시장 기준으로만 조사한다. 국내/한국 시장 수치는 사용하지 않는다.

가장 중요한 규칙:
- 기술명 자체를 검색어/시장명으로 쓰지 말고, 먼저 선정된 대표 적용시장을 기준으로 검색한다.
- 아래의 대표 적용시장 1개를 우선 사용해서 글로벌 시장 CAGR/시장규모/전망을 찾는다.
- display_title은 기술명이 아니라 대표 적용시장 이름으로 작성한다.

대표 적용시장(이미 선정됨):
- 한글 시장명: {rep_market.get('market_name_ko','')}
- 영문 시장명: {rep_market.get('market_name_en','')}
- 권장 표시 제목: {rep_market.get('display_title','')}
- 권장 검색 키워드: {rep_market.get('search_keywords','')}

중요 원칙:
- 웹 검색으로 확인한 실제 공개 출처의 수치만 사용한다.
- CAGR, 시장규모, 기준연도, 전망연도는 출처에 근거해야 한다.
- source_title, source_url, source_year는 가능한 한 정확히 채운다.
- 임의 추정, 보수적 추정, 일반론, 가짜 URL 금지.
- 수치가 충분히 검증되면 verified=true로 하고 CAGR/연도/시장규모를 채운다.
- 수치 검증이 어렵지만 공개 웹페이지에서 관련 시장 그래프 이미지를 찾으면 graph_image_url에 직접 이미지 URL을 넣고, source_title/source_url/source_year/summary를 함께 채운다. 이 경우 verified=false여도 된다.
- graph_image_url은 실제 접근 가능한 공개 이미지 URL만 넣는다. 없으면 빈 문자열.
- {lang_rule}
- market_name은 대표 적용시장 이름으로 작성하고, display_title은 선택 언어에 맞는 시장명으로 작성한다.
- 출력은 JSON만 반환한다.

JSON 형식:
{{
  "market_scope": "global",
  "market_name": "",
  "display_title": "",
  "base_year": "",
  "end_year": "",
  "base_value": "",
  "end_value": "",
  "unit": "",
  "cagr": "",
  "summary": "",
  "source_title": "",
  "source_url": "",
  "source_year": "",
  "graph_image_url": "",
  "graph_image_caption": "",
  "verified": true
}}

기술명(참고용, 시장명으로 쓰지 말 것): {data.get('marketing_title','')}
기술 한줄요약: {data.get('subtitle','')}
적용분야:
{app_text}
"""
    try:
        res = client.responses.create(
            model=TEXT_MODEL_FIXED,
            input=prompt,
            tools=[{"type": "web_search_preview"}],
            temperature=0.1,
        )
        parsed = normalize_market_info(safe_json_parse(res.output_text))
        parsed["market_scope"] = "global"
        parsed["language"] = lang

        # 후처리: display title은 항상 대표 적용시장 중심으로 정리
        default_market_name = parsed.get("market_name") or rep_market.get("market_name_ko") or rep_market.get("market_name_en")
        if not parsed.get("market_name"):
            parsed["market_name"] = default_market_name
        parsed_title = str(parsed.get("display_title") or "").strip()
        tech_title = str(data.get("marketing_title") or "").strip()
        if (not parsed_title) or (tech_title and tech_title in parsed_title):
            parsed["display_title"] = rep_market.get("display_title") or f"글로벌 {default_market_name} 시장"

        parsed["verified"] = bool(parsed.get("verified")) and is_valid_market_info(parsed)
        # 화면에 표시되는 출처명도 선택 언어로 맞춘다. URL/숫자는 그대로 유지.
        if lang != "en" and parsed.get("source_title"):
            parsed["source_title"] = translate_terms_with_gpt({"source_title": parsed.get("source_title", "")}, lang).get("source_title", parsed.get("source_title", ""))
        if not parsed["verified"] and not has_market_graph_image(parsed):
            parsed["error"] = "검증된 수치 또는 활용 가능한 공개 그래프 이미지를 찾지 못했습니다."
        return parsed
    except Exception as e:
        fallback = empty_market_info(f"시장현황 웹 검색 실패: {e}")
        fallback["language"] = lang
        fallback["market_name"] = rep_market.get("market_name_ko") or rep_market.get("market_name_en")
        fallback["display_title"] = rep_market.get("display_title") or fallback.get("display_title")
        return fallback

def build_market_series(info: Dict[str, Any]) -> tuple[list[int], list[float], float, str]:
    m = normalize_market_info(info)
    if not is_valid_market_info(m):
        return [], [], 0.0, ""
    start = _safe_int(m.get("base_year"), 0)
    end = _safe_int(m.get("end_year"), 0)
    if start <= 0 or end <= start:
        return [], [], 0.0, ""
    cagr = _safe_float(m.get("cagr"), 0.0)
    base_value = _safe_float(m.get("base_value"), 0.0)
    end_value = _safe_float(m.get("end_value"), 0.0)
    unit = m.get("unit") or "지수(base=100)"
    n = max(1, end - start)
    if cagr <= 0 and base_value > 0 and end_value > 0:
        try:
            cagr = ((end_value / base_value) ** (1 / n) - 1) * 100.0
        except Exception:
            cagr = 0.0
    if cagr <= 0:
        return [], [], 0.0, ""
    if base_value <= 0:
        base_value = 100.0
        unit = "지수(base=100)"
    years = list(range(start, end + 1))
    values = [base_value * ((1 + cagr / 100.0) ** (yr - start)) for yr in years]
    if end_value > 0 and values:
        scale = end_value / values[-1]
        values = [v * scale for v in values]
    return years, values, cagr, unit

def format_market_value(value: float, unit: str) -> str:
    unit = str(unit or "").strip()
    if unit.startswith("지수"):
        return f"{value:.0f}"
    if value >= 100:
        return f"{value:,.0f}"
    if value >= 10:
        return f"{value:,.1f}"
    return f"{value:,.2f}"

def generate_market_chart(info: Dict[str, Any], primary=(0,55,135), accent=(0,165,185)) -> Image.Image:
    years, values, cagr, unit = build_market_series(info)
    if not years or not values:
        return make_market_placeholder(primary=primary, message1="검증된 글로벌 CAGR 데이터 없음", message2="출처 확인 후 시장현황을 업데이트하세요")
    plt.figure(figsize=(5.1, 2.55), dpi=180)
    ax = plt.gca()
    ax.plot(years, values, marker='o', linewidth=2.5, color=(primary[0]/255, primary[1]/255, primary[2]/255))
    ax.fill_between(years, values, [min(values)]*len(values), alpha=0.08, color=(accent[0]/255, accent[1]/255, accent[2]/255))
    ax.set_xticks(years)
    ax.tick_params(axis='x', labelsize=7)
    ax.tick_params(axis='y', labelsize=7)
    ax.grid(True, axis='y', alpha=0.22)
    ax.set_xlim(min(years), max(years))
    ax.set_ylabel(unit, fontsize=7)
    ax.set_title(f"Global CAGR {cagr:.1f}%", fontsize=9, pad=6)
    for x, y in [(years[0], values[0]), (years[-1], values[-1])]:
        ax.annotate(format_market_value(y, unit), (x, y), textcoords="offset points", xytext=(0, 6), ha='center', fontsize=7)
    plt.tight_layout()
    bio = BytesIO()
    plt.savefig(bio, format='png', bbox_inches='tight', facecolor='white')
    plt.close()
    bio.seek(0)
    return Image.open(bio).convert('RGB')

def _has_hangul(text: str) -> bool:
    return bool(re.search(r"[가-힣]", str(text or "")))

def market_summary_text(info: Dict[str, Any]) -> str:
    return market_description_text(info)

def market_description_text(info: Dict[str, Any]) -> str:
    """시장현황 설명. 한국어는 개조식/음슴체, 다국어는 짧은 bullet style."""
    m = normalize_market_info(info)
    lang = get_lang_code(m.get("language", "ko"))
    years, values, cagr, unit = build_market_series(m)
    lines = []

    if is_valid_market_info(m) and years:
        base_year = years[0]
        end_year = years[-1]
        if lang == "en":
            lines.append(f"• {cagr:.1f}% CAGR projected for {base_year}–{end_year}")
            if values and unit and not str(unit).startswith("지수"):
                end_value = format_market_value(values[-1], unit)
                lines.append(f"• Market expected to reach {end_value} {unit} by {end_year}")
            else:
                lines.append("• Long-term market expansion expected")
        elif lang == "zh":
            lines.append(f"• {base_year}–{end_year}年预计年均增长率{cagr:.1f}%")
            if values and unit and not str(unit).startswith("지수"):
                end_value = format_market_value(values[-1], unit)
                lines.append(f"• {end_year}年市场规模预计达{end_value} {unit}")
            else:
                lines.append("• 中长期市场扩张预期")
        elif lang == "ja":
            lines.append(f"• {base_year}〜{end_year}年の年平均成長率{cagr:.1f}%を予測")
            if values and unit and not str(unit).startswith("지수"):
                end_value = format_market_value(values[-1], unit)
                lines.append(f"• {end_year}年の市場規模は{end_value} {unit}水準を予測")
            else:
                lines.append("• 中長期的な市場拡大を想定")
        else:
            lines.append(f"• {base_year}~{end_year}년 연평균 {cagr:.1f}% 성장 전망")
            if values and unit and not str(unit).startswith("지수"):
                end_value = format_market_value(values[-1], unit)
                lines.append(f"• {end_year}년 시장규모 {end_value} {unit} 수준 전망")
            else:
                lines.append("• 중장기 시장 확대 지속 예상")
        return "\n".join(lines)

    summary = str(m.get("summary") or "").strip()
    if summary:
        if lang == "ko":
            summary = re.sub(r"[.。]$", "", summary)
            return f"• {summary}\n• 세부 수치 출처 기준 확인 필요"
        if lang == "zh":
            return f"• {summary}\n• 需基于来源确认详细数据"
        if lang == "ja":
            return f"• {summary}\n• 詳細数値は出典基準で確認が必要"
        return f"• {summary}\n• Detailed figures require source verification"

    if has_market_graph_image(m):
        if lang == "en":
            return "• Public market data indicates growth trend\n• Detailed figures require source verification"
        if lang == "zh":
            return "• 公开市场资料显示增长趋势\n• 需基于来源确认详细数据"
        if lang == "ja":
            return "• 公開市場資料で成長傾向を確認\n• 詳細数値は出典基準で確認が必要"
        return "• 공개 시장자료 기준 성장 추이 확인됨\n• 세부 수치 출처 기준 확인 필요"

    if lang == "en":
        return "• Verified global market data required\n• Update after confirming source"
    if lang == "zh":
        return "• 需确认可验证的全球市场数据\n• 确认来源后更新市场现状"
    if lang == "ja":
        return "• 検証可能なグローバル市場データが必要\n• 出典確認後に市場動向を反映"
    return "• 검증 가능한 글로벌 시장 데이터 확인 필요\n• 출처 확보 후 시장현황 반영 필요"

def market_source_text(info: Dict[str, Any]) -> str:
    m = normalize_market_info(info)
    lang = get_lang_code(m.get("language", "ko"))
    source_label = label(lang, "source")
    if not (is_valid_market_info(m) or has_market_graph_image(m)):
        return f"{source_label}: " + ({"en":"Data verification required","zh":"需要确认数据","ja":"データ確認が必要"}.get(lang, "데이터 확인 필요"))
    title = m.get("source_title") or source_label
    year = f" ({m.get('source_year')})" if m.get("source_year") else ""
    return f"{source_label}: {title}{year}"


def sanitize_filename_component(value: str) -> str:
    value = str(value or "").strip() or "미입력"
    value = re.sub(r'[\\/:*?"<>|]+', '_', value)
    value = re.sub(r'\s+', ' ', value).strip()
    value = value.replace(' ', '_')
    value = re.sub(r'_+', '_', value)
    return value[:80] or "미입력"

def build_export_basename(data: Dict[str, Any]) -> str:
    ip = normalize_ip(data.get("ip", {}))
    invention_title = ip.get("title") or data.get("original_title") or data.get("marketing_title")
    parts = [
        data.get("university", ""),
        "SMK",
        invention_title,
        data.get("professor", ""),
        data.get("department", ""),
        ip.get("application_number", ""),
    ]
    return "_".join(sanitize_filename_component(p) for p in parts)

def _font_size(font, fallback=18):
    return int(getattr(font, "size", fallback))

def draw_fitted_wrapped(draw, xy, text, size, bold, fill, max_width, max_height, line_gap=6, min_size=12, max_lines=None):
    """박스 밖으로 글자가 나가지 않도록 폰트 크기를 자동 축소해 출력."""
    x, y = xy
    text = str(text or "")
    for fs in range(size, min_size-1, -1):
        font = load_font(fs, bold)
        lines = wrap_text(draw, text, font, max_width, max_lines=max_lines)
        total_h = len(lines) * fs + max(0, len(lines)-1) * line_gap
        if total_h <= max_height:
            yy = y
            for i, line in enumerate(lines):
                if max_lines and i == max_lines-1 and len(wrap_text(draw, text, font, max_width)) > max_lines:
                    line = line[:-1] + "…"
                draw.text((x, yy), line, font=font, fill=fill)
                yy += fs + line_gap
            return yy
    # 그래도 안 맞으면 최소 폰트 + 줄 수 제한
    font = load_font(min_size, bold)
    max_fit_lines = max(1, int(max_height // (min_size + line_gap)))
    lines = wrap_text(draw, text, font, max_width, max_lines=max_lines or max_fit_lines)
    yy = y
    for i, line in enumerate(lines[:max_fit_lines]):
        if i == max_fit_lines-1:
            line = line[:-1] + "…" if len(line) > 1 else "…"
        draw.text((x, yy), line, font=font, fill=fill)
        yy += min_size + line_gap
    return yy

def draw_fitted_centered_wrapped(draw, box, text, size, bold, fill, line_gap=4, min_size=12, max_lines=2):
    """박스 안에서 자동 축소 + 중앙 정렬로 텍스트를 배치한다."""
    x1, y1, x2, y2 = box
    max_width = x2 - x1
    max_height = y2 - y1
    text = str(text or "")
    for fs in range(size, min_size-1, -1):
        font = load_font(fs, bold)
        lines = wrap_text(draw, text, font, max_width, max_lines=max_lines)
        total_h = len(lines) * fs + max(0, len(lines)-1) * line_gap
        if total_h <= max_height:
            yy = y1 + max(0, (max_height - total_h)//2)
            raw_line_count = len(wrap_text(draw, text, font, max_width))
            for i, line in enumerate(lines):
                if max_lines and i == max_lines-1 and raw_line_count > max_lines:
                    line = line[:-1] + "…" if len(line) > 1 else "…"
                bbox = draw.textbbox((0, 0), line, font=font)
                tw = bbox[2] - bbox[0]
                draw.text((x1 + max(0, (max_width - tw)//2), yy), line, font=font, fill=fill)
                yy += fs + line_gap
            return yy
    font = load_font(min_size, bold)
    lines = wrap_text(draw, text, font, max_width, max_lines=max_lines)
    yy = y1 + 2
    for i, line in enumerate(lines):
        if yy > y2 - min_size:
            break
        if i == len(lines)-1:
            line = line[:-1] + "…" if len(line) > 1 else "…"
        bbox = draw.textbbox((0, 0), line, font=font)
        tw = bbox[2] - bbox[0]
        draw.text((x1 + max(0, (max_width - tw)//2), yy), line, font=font, fill=fill)
        yy += min_size + line_gap
    return yy

def draw_bullets_fit(draw, x, y, items, size, bold, color, max_width, max_height, bullet="›", line_gap=5, item_gap=8, min_size=12, max_lines_per_item=3):
    """불릿 목록이 지정 영역을 넘지 않도록 자동 축소."""
    items = [str(v).strip() for v in (items or []) if str(v).strip()]
    for fs in range(size, min_size-1, -1):
        font = load_font(fs, bold)
        yy = y
        ok = True
        cached = []
        for item in items:
            lines = wrap_text(draw, item, font, max_width-26, max_lines=max_lines_per_item)
            h = len(lines) * fs + max(0, len(lines)-1) * line_gap + item_gap
            if yy + h > y + max_height:
                ok = False
                break
            cached.append(lines)
            yy += h
        if ok:
            yy = y
            for lines in cached:
                draw.text((x, yy), bullet, font=font, fill=color)
                tx = x + 26
                for line in lines:
                    draw.text((tx, yy), line, font=font, fill=color)
                    yy += fs + line_gap
                yy += item_gap
            return yy
    # 최소 폰트에서도 넘치면 가능한 만큼만 출력
    font = load_font(min_size, bold)
    yy = y
    bottom = y + max_height
    for item in items:
        if yy >= bottom - min_size:
            break
        lines = wrap_text(draw, item, font, max_width-26, max_lines=max_lines_per_item)
        draw.text((x, yy), bullet, font=font, fill=color)
        tx = x + 26
        for line in lines:
            if yy > bottom - min_size:
                return yy
            draw.text((tx, yy), line, font=font, fill=color)
            yy += min_size + line_gap
        yy += item_gap
    return yy

def get_visual_palette(university_logo: Image.Image | None) -> Dict[str, Tuple[int, int, int]]:
    """대학 로고 대표색을 전체 SMK accent로 사용하되, 본문 가독성은 검정/회색으로 유지한다.
    PIUM 로고는 원본 색상을 그대로 두고, 주변 박스/라인만 대학색 계열 tint로 맞춘다.
    """
    uni_theme = extract_logo_theme(university_logo)
    uni_primary = uni_theme["primary"]

    # 대학 대표색을 기본 accent로 사용
    primary = uni_primary
    primary_dark = ensure_dark(mix(uni_primary, (0, 0, 0), 0.18))
    secondary = mix(primary, (0, 112, 185), 0.15)
    accent = mix(primary, (0, 170, 190), 0.20)

    # 너무 진하게 칠하지 않고, 같은 계열의 연한 tint로 박스/라인 구성
    pale = mix(primary, (255, 255, 255), 0.90)
    pale2 = mix(primary, (255, 255, 255), 0.94)
    pale3 = mix(primary, (255, 255, 255), 0.97)
    line = mix(primary, (255, 255, 255), 0.70)
    line_soft = mix(primary, (255, 255, 255), 0.82)
    table_header = mix(primary, (245, 247, 250), 0.82)

    return {
        "uni_primary": uni_primary,
        "uni_pale": pale,
        "uni_line": line,
        "pium_blue": primary,          # 기존 변수명 호환: 본문 accent도 대학색 계열 적용
        "tech_blue": secondary,
        "tech_cyan": accent,
        "tech_pale": pale2,
        "tech_pale2": pale3,
        "tech_line": line_soft,
        "table_header": table_header,
        "black": (28, 34, 43),
        "gray": (92, 99, 110),
    }



def _draw_shadowed_card(draw, xyxy, radius=28, fill=(255,255,255), outline=(220,230,242), width=2, shadow=True):
    """PPT 수정본처럼 카드에 은은한 그림자를 주는 렌더링 유틸."""
    x1,y1,x2,y2 = xyxy
    if shadow:
        for off, alpha in [(7, 32), (4, 22)]:
            sh = tuple(max(0, min(255, int(c*0.90))) for c in outline)
            draw.rounded_rectangle((x1+off, y1+off, x2+off, y2+off), radius=radius, fill=mix(sh, (255,255,255), 0.55))
    draw.rounded_rectangle(xyxy, radius=radius, fill=fill, outline=outline, width=width)


def compose_tech_brief(data: Dict[str, Any], rep_img: Image.Image, app_imgs: List[Image.Image], contact: str, university_logo: Image.Image | None = None, pium_logo: Image.Image | None = None, qr_img: Image.Image | None = None, piumlink_logo: Image.Image | None = None) -> Image.Image:
    """시장현황 영역을 포함한 1페이지 PIUM Tech Brief 렌더러."""
    W, H = 1240, 1754
    im = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(im)

    pal = get_visual_palette(university_logo)
    uni_primary = pal["uni_primary"]
    uni_pale = pal["uni_pale"]
    uni_line = pal["uni_line"]
    primary = pal["pium_blue"]
    secondary = pal["tech_blue"]
    accent = pal["tech_cyan"]
    sky = pal["tech_pale"]
    sky2 = pal["tech_pale2"]
    line = pal["tech_line"]
    table_header = pal["table_header"]
    black = pal["black"]
    gray = pal["gray"]
    lang = get_lang_code(data.get("language", "ko"))

    d.rectangle((0, 0, W, H), fill=(255,255,255))

    # 긴 기술명이 2~3줄로 바뀌어도 하단 요약 박스와 겹치지 않도록
    # 헤더 영역을 기존 248px에서 280px로 확장하고, 제목/요약 박스 좌표를 분리한다.
    header_h = 280
    d.rectangle((0, 0, W, header_h), fill=uni_pale)
    d.line((0, header_h, W, header_h), fill=mix(line, (160,160,160), 0.15), width=2)

    uni_logo_size = 124
    uni_logo_x, uni_logo_y = 28, 24
    left_logo = make_transparent_logo_canvas(university_logo, size=(uni_logo_size, uni_logo_size), padding=0) if university_logo else make_university_logo_box(None, data.get('university_display') or data.get('university',''), size=(uni_logo_size, uni_logo_size), bg=uni_pale, primary=uni_primary)
    im.paste(left_logo, (uni_logo_x, uni_logo_y), left_logo if left_logo.mode == 'RGBA' else None)

    # PIUM 로고/QR은 원본 컬러를 보존해야 하므로, 대학색 헤더와 직접 섞이지 않게
    # 흰색 독립 카드 안에 배치한다. 어떤 대학 대표색이 와도 브랜드 충돌을 줄이는 보호 영역.
    right_card_x, right_card_y = W - 178, 18
    right_card_w, right_card_h = 148, 220
    _draw_shadowed_card(
        d,
        (right_card_x, right_card_y, right_card_x + right_card_w, right_card_y + right_card_h),
        radius=18,
        fill=(255, 255, 255),
        outline=uni_line,
        width=1,
        shadow=True,
    )
    right_x = right_card_x + 10
    top_y = right_card_y + 14
    if pium_logo is not None:
        pium_canvas = make_transparent_logo_canvas(pium_logo, size=(128, 48), padding=0)
        im.paste(pium_canvas, (right_x, top_y), pium_canvas)
    if piumlink_logo is not None:
        link_size = 112
        link = make_transparent_logo_canvas(piumlink_logo, size=(link_size, link_size), padding=0)
        im.paste(link, (right_card_x + (right_card_w - link_size)//2, top_y + 60), link)

    header_x = 190
    header_w = right_card_x - header_x - 24
    prof_suffix = label(lang, "prof_suffix")
    kicker = f"PIUM Tech Offer  x  {data.get('university_display') or data.get('university','')}  |  {data.get('department_display') or data.get('department','')}  |  {data.get('professor','')} {prof_suffix}"
    draw_fitted_wrapped(d, (header_x, 42), kicker, 25, False, uni_primary, header_w, 32, line_gap=2, min_size=17, max_lines=1)
    draw_fitted_wrapped(
        d,
        (header_x, 78),
        data.get("marketing_title", "기술명"),
        41,
        True,
        uni_primary,
        header_w,
        108,
        line_gap=4,
        min_size=24,
        max_lines=3,
    )

    # 서브카피는 별도 박스 없이 제목 아래에 자연스럽게 이어지는 한 줄 설명으로 처리
    subtitle = (data.get("subtitle", "") or "").strip()
    if subtitle:
        sub_text = ensure_dark(mix(uni_primary, black, 0.35))
        draw_fitted_wrapped(
            d,
            (header_x, 202),
            f"- {subtitle}",
            19,
            False,
            sub_text,
            header_w,
            34,
            line_gap=2,
            min_size=12,
            max_lines=2,
        )

    X = 28
    CW = W - 56
    card_line = line
    sec_font = load_font(26, True)

    # Applications + Market
    app_y, app_h = 294, 262
    top_gap = 30
    app_card_w = 552
    market_x = X + app_card_w + top_gap
    market_w = CW - app_card_w - top_gap

    # 참고용 레이아웃처럼 적용분야와 시장현황을 완전히 독립된 카드로 분리
    _draw_shadowed_card(d, (X, app_y, X+app_card_w, app_y+app_h), radius=34, fill=(255,255,255), outline=card_line, width=2, shadow=True)
    _draw_shadowed_card(d, (market_x, app_y, market_x+market_w, app_y+app_h), radius=34, fill=(255,255,255), outline=card_line, width=2, shadow=True)
    draw_section_title(d, X+32, app_y+28, label(lang, "apps"), sec_font, primary)
    draw_section_title(d, market_x+34, app_y+28, label(lang, "market"), sec_font, primary)

    apps = data.get("applications", [])[:3]
    app_inner_x = X + 22
    app_inner_w = app_card_w - 44
    col_gap = 4
    col_w = (app_inner_w - col_gap*2) // 3
    icon_size = (112, 88)
    icon_y = app_y + 82
    for i in range(3):
        col_x = app_inner_x + i*(col_w + col_gap)
        if i > 0:
            sep_x = col_x - col_gap//2
            d.line((sep_x, app_y+84, sep_x, app_y+204), fill=mix(card_line, (255,255,255), 0.24), width=1)
        if i < len(app_imgs):
            icon = fit_image(app_imgs[i], icon_size, bg=(255,255,255), trim=True)
            im.paste(icon, (col_x + (col_w-icon_size[0])//2, icon_y))
        app = apps[i] if i < len(apps) else {"name":""}
        draw_centered_wrapped(d, (col_x+6, app_y+176, col_x+col_w-6, app_y+228), app.get("name", ""), load_font(19, True), black, max_lines=2, line_gap=3)

    market_info = normalize_market_info(data.get("market_info", {}))
    market_title = market_info.get("display_title") or market_info.get("market_name") or label(lang, "market_default")
    draw_fitted_wrapped(d, (market_x+34, app_y+74), market_title, 19, True, black, market_w-68, 30, line_gap=2, min_size=14, max_lines=1)

    # 시장현황 카드 내부: 왼쪽 그래프 + 오른쪽 설명 카드 + 하단 출처
    chart, _market_visual_mode = get_market_visual(market_info, primary=primary, accent=accent)
    graph_x, graph_y = market_x+34, app_y+104
    graph_w, graph_h = 178, 104
    _draw_shadowed_card(d, (graph_x, graph_y, graph_x+graph_w, graph_y+graph_h), radius=12, fill=(255,255,255), outline=card_line, width=1, shadow=False)
    chart_img = fit_image(chart, (graph_w-10, graph_h-10), bg=(255,255,255), trim=False)
    im.paste(chart_img, (graph_x+5, graph_y+5))

    desc_x = graph_x + graph_w + 16
    desc_y = app_y + 114
    desc_w = market_x + market_w - 28 - desc_x
    desc_h = 78
    _draw_shadowed_card(d, (desc_x, desc_y, desc_x+desc_w, desc_y+desc_h), radius=16, fill=(255,255,255), outline=card_line, width=1, shadow=True)
    desc = market_description_text(market_info)
    draw_fitted_wrapped(d, (desc_x+14, desc_y+12), desc, 13, False, gray, desc_w-28, desc_h-18, line_gap=3, min_size=10, max_lines=4)

    source = market_source_text(market_info)
    draw_fitted_wrapped(d, (graph_x, app_y+218), source, 11, False, gray, market_w-64, 18, line_gap=2, min_size=8, max_lines=1)

    # Overview / Differentiation
    y = 582
    gap = 30
    box_w = (CW - gap) // 2
    left_x = X
    right_x2 = X + box_w + gap
    box_h = 290
    for bx, title, items in [
        (left_x, label(lang, "overview"), data.get("overview", [])[:3]),
        (right_x2, label(lang, "diff"), data.get("differentiation", [])[:3]),
    ]:
        _draw_shadowed_card(d, (bx, y, bx+box_w, y+box_h), radius=28, fill=sky, outline=card_line, width=2, shadow=True)
        draw_section_title(d, bx+36, y+34, title, sec_font, primary)
        draw_bullets_fit(d, bx+42, y+88, items, 19, False, black, box_w-84, box_h-112, bullet="›", line_gap=4, item_gap=8, min_size=13, max_lines_per_item=3)

    # Representative / Competitiveness
    y = 900
    rep_w_box = 380
    comp_x = X + rep_w_box + 22
    comp_w = X + CW - comp_x
    comp_h = 322
    _draw_shadowed_card(d, (X, y, X+rep_w_box, y+comp_h), radius=28, fill=(255,255,255), outline=card_line, width=2, shadow=True)
    _draw_shadowed_card(d, (comp_x, y, comp_x+comp_w, y+comp_h), radius=28, fill=(255,255,255), outline=card_line, width=2, shadow=True)
    draw_section_title(d, X+36, y+34, label(lang, "drawing"), sec_font, primary)
    draw_section_title(d, comp_x+36, y+34, label(lang, "competitiveness"), sec_font, primary)

    rep = fit_image(rep_img, (rep_w_box-74, 214), bg=(255,255,255), trim=True)
    im.paste(rep, (X+37, y+82))

    inner_x = comp_x + 44
    inner_w = comp_w - 70
    sub1_y, sub1_h = y+76, 102
    sub2_y, sub2_h = y+194, 108
    _draw_shadowed_card(d, (inner_x, sub1_y, inner_x+inner_w, sub1_y+sub1_h), radius=16, fill=(255,255,255), outline=card_line, width=1, shadow=True)
    d.text((inner_x+20, sub1_y+14), label(lang, "limitations"), font=load_font(19, True), fill=gray)
    draw_bullets_fit(d, inner_x+22, sub1_y+42, data.get("limitations", [])[:2], 14, False, black, inner_w-44, sub1_h-48, bullet="•", line_gap=3, item_gap=2, min_size=11, max_lines_per_item=2)

    _draw_shadowed_card(d, (inner_x, sub2_y, inner_x+inner_w, sub2_y+sub2_h), radius=16, fill=sky2, outline=card_line, width=1, shadow=True)
    d.text((inner_x+20, sub2_y+14), label(lang, "advantages"), font=load_font(19, True), fill=secondary)
    draw_bullets_fit(d, inner_x+22, sub2_y+42, data.get("technical_advantages", [])[:2], 14, False, primary, inner_w-44, sub2_h-48, bullet="•", line_gap=3, item_gap=2, min_size=11, max_lines_per_item=2)

    # IP full width
    y = 1252
    ip = normalize_ip(data.get("ip", {}))
    ip_h = 220
    _draw_shadowed_card(d, (X, y, X+CW, y+ip_h), radius=28, fill=sky2, outline=card_line, width=2, shadow=True)
    draw_section_title(d, X+36, y+30, label(lang, "ip"), sec_font, primary)

    table_x, table_y = X+28, y+72
    table_w, table_h = CW-56, 132
    header_h2 = 40
    col1 = int(table_w*0.45)
    col2 = int(table_w*0.28)
    col3 = table_w - col1 - col2
    d.rectangle((table_x, table_y, table_x+table_w, table_y+header_h2), fill=table_header, outline=card_line)
    d.rectangle((table_x, table_y+header_h2, table_x+table_w, table_y+table_h), fill=(255,255,255), outline=card_line)
    c1 = table_x + col1
    c2 = c1 + col2
    for xx in [c1, c2]:
        d.line((xx, table_y, xx, table_y+table_h), fill=card_line, width=1)
    draw_centered_wrapped(d, (table_x+8, table_y+4, c1-8, table_y+header_h2-2), label(lang, "invention"), load_font(19, True), black, max_lines=1)
    draw_centered_wrapped(d, (c1+8, table_y+3, c2-8, table_y+header_h2-2), ip_table_label(lang, "app_no", ip), load_font(14, True), black, max_lines=2, line_gap=1)
    draw_centered_wrapped(d, (c2+8, table_y+3, table_x+table_w-8, table_y+header_h2-2), ip_table_label(lang, "app_date", ip), load_font(14, True), black, max_lines=2, line_gap=1)
    draw_centered_wrapped(d, (table_x+12, table_y+header_h2+8, c1-12, table_y+table_h-8), ip["title"] or data.get("original_title", ""), load_font(15, False), black, max_lines=2, line_gap=2)
    num_text = f"{ip['application_number']}\n({ip['registration_number']})" if ip['registration_number'] else ip['application_number']
    date_text = f"{ip['application_date']}\n({ip['registration_date']})" if ip['registration_date'] else ip['application_date']
    draw_centered_wrapped(d, (c1+10, table_y+header_h2+8, c2-10, table_y+table_h-8), num_text, load_font(17, False), black, max_lines=2, line_gap=3)
    draw_centered_wrapped(d, (c2+10, table_y+header_h2+8, table_x+table_w-10, table_y+table_h-8), date_text, load_font(17, False), black, max_lines=2, line_gap=3)

    # Contact
    y = 1504
    contact_h = 78
    _draw_shadowed_card(d, (X, y, X+CW, y+contact_h), radius=18, fill=(255,255,255), outline=card_line, width=2, shadow=True)
    d.text((X+42, y+23), label(lang, "contact"), font=load_font(26, True), fill=primary)
    draw_fitted_wrapped(d, (X+160, y+27), contact, 20, False, black, CW-190, 30, line_gap=3, min_size=13, max_lines=1)

    return im

# -----------------------------------------------------
# PDF / PPTX
# -----------------------------------------------------
def _image_to_a4_pdf_bytes(img: Image.Image, dpi: int = 300) -> bytes:
    """이미지 1장을 실제 A4 페이지 크기의 고해상도 PDF로 저장한다.

    PIL의 img.save(..., "PDF", resolution=300)는 이미지 픽셀/dpi 값으로
    PDF 페이지 크기를 계산하기 때문에, 1024x1536 이미지는 약 3.4x5.1 inch의
    작은 PDF 페이지가 되어 뷰어에서 확대 시 흐려 보일 수 있다.
    이 함수는 PDF 페이지를 A4 포인트 크기로 고정하고, 이미지를 A4 안에
    300dpi급 픽셀로 리샘플링한 뒤 PNG 스트림으로 삽입한다.
    """
    a4_w_pt, a4_h_pt = 595.2756, 841.8898  # A4 in PDF points
    src = img.convert("RGB")
    src_ratio = src.width / max(1, src.height)
    page_ratio = a4_w_pt / a4_h_pt

    if src_ratio >= page_ratio:
        draw_w_pt = a4_w_pt
        draw_h_pt = a4_w_pt / src_ratio
    else:
        draw_h_pt = a4_h_pt
        draw_w_pt = a4_h_pt * src_ratio

    x0 = (a4_w_pt - draw_w_pt) / 2
    y0 = (a4_h_pt - draw_h_pt) / 2
    rect = fitz.Rect(x0, y0, x0 + draw_w_pt, y0 + draw_h_pt)

    target_w_px = max(1, int(round(draw_w_pt / 72 * dpi)))
    target_h_px = max(1, int(round(draw_h_pt / 72 * dpi)))
    resampled = src.resize((target_w_px, target_h_px), Image.LANCZOS)
    resampled = ImageEnhance.Sharpness(resampled).enhance(1.04)

    png = BytesIO()
    resampled.save(png, format="PNG", optimize=True)
    png_bytes = png.getvalue()

    doc = fitz.open()
    page = doc.new_page(width=a4_w_pt, height=a4_h_pt)
    page.insert_image(rect, stream=png_bytes, keep_proportion=True)
    out = doc.tobytes(deflate=True, garbage=4)
    doc.close()
    return out


def make_pdf_bytes_from_image(img: Image.Image) -> bytes:
    # 기본 SMK PDF도 실제 A4 페이지로 저장해 PDF 뷰어 확대 시 품질 저하를 줄인다.
    return _image_to_a4_pdf_bytes(img, dpi=300)




def generate_market_bar_chart(info: Dict[str, Any], primary=(0,55,135), accent=(0,165,185)) -> Image.Image:
    years, values, cagr, unit = build_market_series(info)
    if not years or not values:
        return make_market_placeholder(primary=primary, message1="검증된 글로벌 CAGR 데이터 없음", message2="출처 확인 후 시장현황을 업데이트하세요")
    plt.figure(figsize=(5.2, 2.85), dpi=180)
    ax = plt.gca()
    colors = []
    total = max(1, len(years)-1)
    for i in range(len(years)):
        ratio = i / total
        rr = mix(primary, (255,255,255), max(0.16, 0.50 - ratio*0.30))
        colors.append((rr[0]/255, rr[1]/255, rr[2]/255))
    ax.bar(years, values, color=colors, width=0.62)
    ax.grid(True, axis='y', alpha=0.18)
    ax.tick_params(axis='x', labelsize=7)
    ax.tick_params(axis='y', labelsize=7)
    ax.set_xticks(years)
    ax.set_xlim(min(years)-0.6, max(years)+0.6)
    ax.set_ylabel(unit, fontsize=7)
    ax.set_title('Annual Market Growth', fontsize=9, pad=6)
    # annotate first / middle / last for readability
    idxs = {0, len(years)-1, len(years)//2}
    for idx in sorted(idxs):
        x, y = years[idx], values[idx]
        ax.annotate(format_market_value(y, unit), (x, y), textcoords='offset points', xytext=(0, 4), ha='center', fontsize=6.8)
    for spine in ['top', 'right']:
        ax.spines[spine].set_visible(False)
    plt.tight_layout()
    bio = BytesIO()
    plt.savefig(bio, format='png', bbox_inches='tight', facecolor='white')
    plt.close()
    bio.seek(0)
    return Image.open(bio).convert('RGB')

HIGHLIGHT_TERMS = [
    '실시간 모니터링', '데이터 진위성', '데이터 무결성', '안전하게 저장', '블록체인', '스마트 컨트랙트',
    '엣지 노드', '유무선 통신', '무선 및 유선 통신', '네트워크 환경', '호환 가능', '독립적으로 운영',
    '확장성 제공', '운영 모니터링', '자동 추출 및 관리', '품질 이력 관리', '공급망 투명성', '신뢰성 확보',
    '투명성 강화', '실시간 개입', '분석 및 해설', '실시간 분석', '고객지원', '온라인 교육', '실시간 대응',
    '품질 보증', '경제적 이득', '위변조 방지', '상태 실시간 모니터링', '생산 설비 관리', '공정 관리'
]


def _highlight_spans(text: str):
    src = str(text or '').strip()
    spans = []
    for term in sorted(HIGHLIGHT_TERMS, key=len, reverse=True):
        start = 0
        while True:
            idx = src.find(term, start)
            if idx < 0:
                break
            span = (idx, idx + len(term))
            if not any(max(a, span[0]) < min(b, span[1]) for a, b in spans):
                spans.append(span)
            start = idx + len(term)
    if not spans and src:
        tokens = src.split()
        if len(tokens) >= 2:
            phrase = ' '.join(tokens[:2])
            spans.append((0, len(phrase)))
        elif len(src) > 8:
            spans.append((0, min(len(src), 10)))
    return sorted(spans)


def segment_text_with_emphasis(text: str):
    src = str(text or '').strip()
    spans = _highlight_spans(src)
    if not spans:
        return [(src, False)] if src else []
    segs = []
    cur = 0
    for a, b in spans:
        if a > cur:
            segs.append((src[cur:a], False))
        segs.append((src[a:b], True))
        cur = b
    if cur < len(src):
        segs.append((src[cur:], False))
    return [(t, b) for t, b in segs if t]


def _tokenize_rich_segments(segments):
    out = []
    for text, bold in segments:
        parts = re.findall(r'\S+|\s+', text)
        for p in parts:
            out.append((p, bold))
    return out


def _text_width(draw, text, font):
    bb = draw.textbbox((0, 0), text, font=font)
    return bb[2] - bb[0]


def _wrap_rich_segments(draw, segments, font_regular, font_bold, max_width, max_lines=None):
    tokens = _tokenize_rich_segments(segments)
    lines = []
    cur = []
    cur_w = 0
    for tok, is_bold in tokens:
        font = font_bold if is_bold else font_regular
        if tok.isspace() and not cur:
            continue
        tok_w = _text_width(draw, tok, font)
        if not tok.isspace() and tok_w > max_width:
            chunks = _split_token_to_fit(draw, tok, font, max_width)
            for j, ch in enumerate(chunks):
                ch_w = _text_width(draw, ch, font)
                if cur and cur_w + ch_w > max_width:
                    while cur and cur[-1][0].isspace():
                        cur.pop()
                    if cur:
                        lines.append(cur)
                        if max_lines and len(lines) >= max_lines:
                            return lines
                    cur, cur_w = [], 0
                cur.append((ch, is_bold))
                cur_w += ch_w
                if j < len(chunks)-1:
                    while cur and cur[-1][0].isspace():
                        cur.pop()
                    lines.append(cur)
                    if max_lines and len(lines) >= max_lines:
                        return lines
                    cur, cur_w = [], 0
            continue
        if cur and cur_w + tok_w > max_width:
            while cur and cur[-1][0].isspace():
                cur.pop()
            lines.append(cur)
            if max_lines and len(lines) >= max_lines:
                return lines
            cur = []
            cur_w = 0
            tok = tok.lstrip()
            if not tok:
                continue
            tok_w = _text_width(draw, tok, font)
        cur.append((tok, is_bold))
        cur_w += tok_w
    while cur and cur[-1][0].isspace():
        cur.pop()
    if cur:
        lines.append(cur)
    if max_lines:
        lines = lines[:max_lines]
    return lines


def draw_emphasis_bullets_fit(draw, x, y, items, size, color, highlight_color, max_width, max_height, bullet='•', line_gap=5, item_gap=8, min_size=12, max_lines_per_item=3):
    items = [str(v).strip() for v in (items or []) if str(v).strip()]
    for fs in range(size, min_size-1, -1):
        fr = load_font(fs, False)
        fb = load_font(fs, True)
        yy = y
        layouts = []
        ok = True
        for item in items:
            segs = segment_text_with_emphasis(item)
            lines = _wrap_rich_segments(draw, segs, fr, fb, max_width-28, max_lines=max_lines_per_item)
            hh = len(lines) * fs + max(0, len(lines)-1) * line_gap + item_gap
            if yy + hh > y + max_height:
                ok = False
                break
            layouts.append(lines)
            yy += hh
        if ok:
            yy = y
            for lines in layouts:
                draw.text((x, yy), bullet, font=fr, fill=color)
                for line in lines:
                    tx = x + 28
                    for tok, is_bold in line:
                        font = fb if is_bold else fr
                        fill = highlight_color if is_bold else color
                        draw.text((tx, yy), tok, font=font, fill=fill)
                        tx += _text_width(draw, tok, font)
                    yy += fs + line_gap
                yy += item_gap
            return yy
    # minimum fallback
    fr = load_font(min_size, False)
    fb = load_font(min_size, True)
    yy = y
    bottom = y + max_height
    for item in items:
        if yy > bottom - min_size:
            break
        segs = segment_text_with_emphasis(item)
        lines = _wrap_rich_segments(draw, segs, fr, fb, max_width-28, max_lines=max_lines_per_item)
        draw.text((x, yy), bullet, font=fr, fill=color)
        for line in lines:
            if yy > bottom - min_size:
                break
            tx = x + 28
            for tok, is_bold in line:
                font = fb if is_bold else fr
                fill = highlight_color if is_bold else color
                draw.text((tx, yy), tok, font=font, fill=fill)
                tx += _text_width(draw, tok, font)
            yy += min_size + line_gap
        yy += item_gap
    return yy



def draw_linear_gradient(im: Image.Image, box, color1, color2, horizontal=False):
    x1, y1, x2, y2 = map(int, box)
    w = max(1, x2 - x1)
    h = max(1, y2 - y1)
    grad = Image.new('RGB', (w, h), color1)
    px = grad.load()
    if horizontal:
        denom = max(1, w - 1)
        for x in range(w):
            t = x / denom
            c = tuple(int(color1[i] * (1 - t) + color2[i] * t) for i in range(3))
            for y in range(h):
                px[x, y] = c
    else:
        denom = max(1, h - 1)
        for y in range(h):
            t = y / denom
            c = tuple(int(color1[i] * (1 - t) + color2[i] * t) for i in range(3))
            for x in range(w):
                px[x, y] = c
    im.paste(grad, (x1, y1))


def draw_network_pattern(draw: ImageDraw.ImageDraw, box, color=(220,230,242), width=2):
    x1, y1, x2, y2 = box
    pts = [
        (x1 + 12, y2 - 12), (x1 + 46, y2 - 42), (x1 + 92, y2 - 24),
        (x1 + 130, y2 - 58), (x1 + 176, y2 - 40), (x1 + 220, y2 - 72),
        (x2 - 180, y1 + 38), (x2 - 138, y1 + 18), (x2 - 98, y1 + 56),
        (x2 - 58, y1 + 34), (x2 - 24, y1 + 78),
    ]
    links = [(0,1),(1,2),(2,3),(3,4),(6,7),(7,8),(8,9),(9,10),(2,5)]
    for a,b in links:
        if a < len(pts) and b < len(pts):
            draw.line((pts[a][0], pts[a][1], pts[b][0], pts[b][1]), fill=color, width=width)
    for x, y in pts:
        r = 4
        draw.ellipse((x-r, y-r, x+r, y+r), fill=color, outline=color)


def get_premium_app_desc(name: str, idx: int) -> str:
    name = str(name or '')
    presets = [
        ('스마트 팩토리', '실시간 데이터 수집·분석으로\n스마트 공정 최적화 실현'),
        ('생산 설비 관리', '설비 상태 모니터링 및\n예지정비로 가동률 극대화'),
        ('공급망 투명성', '블록체인 기반 이력관리로\n공급망 신뢰 및 추적성 강화'),
        ('온라인 교육 플랫폼', '실시간 상호작용 기반\n학습 품질 및 몰입도 향상'),
        ('고객 상담 서비스', '상담 대화 분석과 개입으로\n이해도 및 응대 품질 향상'),
    ]
    for key, desc in presets:
        if key in name:
            return desc
    return presets[idx % len(presets)][1]


def draw_simple_feature_icon(draw: ImageDraw.ImageDraw, box, kind: str, color=(255,255,255)):
    x1, y1, x2, y2 = box
    cx = (x1 + x2) // 2
    cy = (y1 + y2) // 2
    if kind == 'monitor':
        draw.rounded_rectangle((x1+6, y1+8, x2-6, y2-16), radius=6, outline=color, width=3)
        draw.line((cx, y2-16, cx, y2-8), fill=color, width=3)
        draw.line((cx-12, y2-6, cx+12, y2-6), fill=color, width=3)
    elif kind == 'database':
        draw.ellipse((x1+10, y1+8, x2-10, y1+20), outline=color, width=3)
        draw.rectangle((x1+10, y1+14, x2-10, y2-14), outline=color, width=3)
        draw.arc((x1+10, y2-22, x2-10, y2-10), 0, 180, fill=color, width=3)
    elif kind == 'network':
        pts = [(cx, y1+10), (x1+14, cy+2), (x2-14, cy+2), (cx, y2-12)]
        for a,b in [(0,1),(0,2),(1,3),(2,3)]:
            draw.line((pts[a][0], pts[a][1], pts[b][0], pts[b][1]), fill=color, width=3)
        for x,y in pts:
            draw.ellipse((x-4,y-4,x+4,y+4), fill=color, outline=color)
    elif kind == 'shield':
        draw.polygon([(cx,y1+8),(x2-12,y1+18),(x2-16,y2-14),(cx,y2-8),(x1+16,y2-14),(x1+12,y1+18)], outline=color, fill=None, width=3)
        draw.line((cx-8, cy, cx-1, cy+7), fill=color, width=3)
        draw.line((cx-1, cy+7, cx+10, cy-6), fill=color, width=3)
    else:
        draw.ellipse((x1+10,y1+10,x2-10,y2-10), outline=color, width=3)


def app_icon_kind(name: str, idx: int) -> str:
    name = str(name or '')
    if '스마트 팩토리' in name or '메신저' in name or '교육' in name:
        return 'monitor'
    if '설비' in name or '관리' in name:
        return 'database'
    if '공급망' in name or '투명' in name:
        return 'network'
    if '고객' in name or '상담' in name:
        return 'shield'
    return ['monitor', 'database', 'network'][idx % 3]


def draw_arrow_down(im: Image.Image, center_x: int, y: int, color=(160,200,245)):
    d = ImageDraw.Draw(im)
    pts = [(center_x-34, y), (center_x+34, y), (center_x, y+28)]
    d.polygon(pts, fill=color)
    d.polygon([(center_x-24, y-10), (center_x+24, y-10), (center_x, y+10)], fill=mix(color, (255,255,255), 0.15))

def compose_infographic_brief(
    data: Dict[str, Any],
    rep_img: Image.Image,
    app_imgs: List[Image.Image],
    contact: str,
    university_logo: Image.Image | None = None,
    pium_logo: Image.Image | None = None,
    qr_img: Image.Image | None = None,
    piumlink_logo: Image.Image | None = None,
) -> Image.Image:
    """예시 시안과 유사한 브로슈어형 프리미엄 인포그래픽."""
    W, H = 1240, 1754
    im = Image.new('RGB', (W, H), (248, 250, 253))
    d = ImageDraw.Draw(im)

    pal = get_visual_palette(university_logo)
    primary = ensure_dark(pal['uni_primary'])
    secondary = ensure_dark(mix(primary, (47, 120, 214), 0.25))
    accent = mix(primary, (0, 188, 212), 0.30)
    black = (32, 40, 52)
    gray = (90, 98, 112)
    light_bg = (248, 250, 253)
    line = mix(primary, (214,224,238), 0.80)
    deep_blue = ensure_dark(mix(primary, (0, 70, 170), 0.15))
    green_tint = mix((155, 225, 205), (255,255,255), 0.55)
    blue_tint = mix((132, 186, 248), (255,255,255), 0.55)
    lang = get_lang_code(data.get('language', 'ko'))

    d.rectangle((0,0,W,H), fill=light_bg)
    draw_network_pattern(d, (20, 120, 340, 330), color=mix(primary, (255,255,255), 0.86), width=2)
    draw_network_pattern(d, (920, 80, 1210, 250), color=mix(primary, (255,255,255), 0.86), width=2)
    draw_network_pattern(d, (1000, 1650, 1220, 1736), color=mix(primary, (255,255,255), 0.82), width=2)

    if university_logo is not None:
        logo = make_transparent_logo_canvas(university_logo, size=(118,118), padding=0)
        im.paste(logo, (28, 24), logo)
    else:
        fallback = make_university_logo_box(None, data.get('university_display') or data.get('university',''), size=(118,118), bg=(255,255,255), primary=primary)
        im.paste(fallback, (28, 24), fallback if fallback.mode == 'RGBA' else None)

    qr_card = (1016, 20, 1210, 264)
    _draw_shadowed_card(d, qr_card, radius=24, fill=(255,255,255), outline=line, width=2, shadow=True)
    if pium_logo is not None:
        pl = make_transparent_logo_canvas(pium_logo, size=(124, 48), padding=0)
        im.paste(pl, (1052, 36), pl)
    qr_source = qr_img or piumlink_logo
    if qr_source is not None:
        qr_fit = fit_image(qr_source.convert('RGB') if qr_source.mode != 'RGB' else qr_source, (134, 134), bg=(255,255,255), trim=True)
        qr_canvas = Image.new('RGB', (134, 134), 'white')
        qr_canvas.paste(qr_fit, ((134-qr_fit.width)//2, (134-qr_fit.height)//2))
        im.paste(qr_canvas, (1047, 92))

    head_x = 186
    head_w = 794
    prof_suffix = label(lang, 'prof_suffix')
    kicker = f"PIUM Tech Offer x {data.get('university_display') or data.get('university','')} | {data.get('department_display') or data.get('department','')} | {data.get('professor','')} {prof_suffix}"
    draw_fitted_wrapped(d, (head_x, 40), kicker, 16, False, gray, head_w, 20, line_gap=2, min_size=11, max_lines=1)
    draw_fitted_wrapped(d, (head_x, 78), data.get('marketing_title', '기술명'), 42, True, primary, head_w, 108, line_gap=4, min_size=24, max_lines=3)
    subtitle = str(data.get('subtitle') or '').strip()
    if subtitle:
        draw_fitted_wrapped(d, (head_x, 192), subtitle, 18, False, black, head_w-30, 22, line_gap=2, min_size=12, max_lines=2)

    top_y = 280
    left_box = (24, top_y, 654, 670)
    right_box = (670, top_y, 1216, 670)
    for box in [left_box, right_box]:
        _draw_shadowed_card(d, box, radius=24, fill=(255,255,255), outline=line, width=2, shadow=True)

    d.rounded_rectangle((38, top_y+14, 174, top_y+58), radius=20, fill=deep_blue)
    d.text((58, top_y+22), label(lang, 'apps'), font=load_font(22, True), fill=(255,255,255))
    d.rounded_rectangle((690, top_y+14, 814, top_y+58), radius=20, fill=deep_blue)
    d.text((712, top_y+22), label(lang, 'market'), font=load_font(22, True), fill=(255,255,255))

    app_names = [a.get('name','') for a in data.get('applications', [])[:3]]
    while len(app_names) < 3:
        app_names.append('')
    app_colors = [((18, 92, 200), (31, 146, 240)), ((19, 118, 142), (38, 164, 173)), ((74, 59, 174), (132, 66, 204))]
    card_w = 184
    card_h = 300
    start_x = 38
    card_top = top_y + 76
    for i in range(3):
        x = start_x + i*(card_w + 18)
        box = (x, card_top, x+card_w, card_top+card_h)
        _draw_shadowed_card(d, box, radius=22, fill=(255,255,255), outline=line, width=1, shadow=False)
        inner = (x+4, card_top+4, x+card_w-4, card_top+card_h-4)
        draw_linear_gradient(im, inner, app_colors[i][0], app_colors[i][1], horizontal=False)
        dd = ImageDraw.Draw(im)
        dd.rounded_rectangle(inner, radius=20, outline=None)
        dd.rounded_rectangle((x+10, card_top+10, x+card_w-10, card_top+48), radius=14, fill=mix((255,255,255), app_colors[i][1], 0.80))
        name = app_names[i]
        desc = get_premium_app_desc(name, i)
        draw_fitted_centered_wrapped(dd, (x+12, card_top+16, x+card_w-12, card_top+42), name, 18, True, (255,255,255), line_gap=1, min_size=11, max_lines=1)
        dd.rounded_rectangle((x+18, card_top+64, x+card_w-18, card_top+210), radius=18, fill=mix(app_colors[i][0], (255,255,255), 0.20))
        if i < len(app_imgs):
            icon = fit_image(app_imgs[i], (120, 112), bg=mix(app_colors[i][0], (255,255,255), 0.20), trim=True)
            im.paste(icon, (x + (card_w-icon.width)//2, card_top+78))
        kind = app_icon_kind(name, i)
        draw_simple_feature_icon(dd, (x+18, card_top+218, x+52, card_top+252), kind, color=(255,255,255))
        draw_fitted_centered_wrapped(dd, (x+20, card_top+224, x+card_w-20, card_top+286), desc, 11, False, (255,255,255), line_gap=2, min_size=9, max_lines=3)

    market_info = normalize_market_info(data.get('market_info', {}))
    market_title = market_info.get('display_title') or market_info.get('market_name') or label(lang, 'market_default')
    draw_fitted_wrapped(d, (694, top_y+78), market_title, 20, True, black, 420, 24, line_gap=2, min_size=12, max_lines=2)
    bar_chart = generate_market_bar_chart(market_info, primary=deep_blue, accent=accent)
    chart_box = (694, top_y+110, 1000, top_y+302)
    d.rounded_rectangle(chart_box, radius=10, fill=(255,255,255))
    chart_img = fit_image(bar_chart, (chart_box[2]-chart_box[0], chart_box[3]-chart_box[1]), bg=(255,255,255), trim=False)
    im.paste(chart_img, (chart_box[0], chart_box[1]))
    years, values, cagr, unit = build_market_series(market_info)
    end_year = years[-1] if years else _safe_int(market_info.get('end_year'), 0)
    end_value = format_market_value(values[-1], unit) if values else str(market_info.get('end_value') or '-')
    k1 = (1024, top_y+128, 1188, top_y+224)
    k2 = (1024, top_y+244, 1188, top_y+340)
    for box, icon_kind, headline, body in [
        (k1, 'trend', f'2025~{end_year}년', f'연평균 {cagr:.1f}%\n성장 전망'),
        (k2, 'globe', f'{end_year}년 시장규모', f'{end_value}\n{unit} 수준 전망'),
    ]:
        _draw_shadowed_card(d, box, radius=16, fill=(255,255,255), outline=line, width=1, shadow=False)
        d.rounded_rectangle((box[0]+12, box[1]+18, box[0]+52, box[1]+58), radius=10, fill=mix(primary, (255,255,255), 0.88))
        if icon_kind == 'trend':
            d.line((box[0]+18, box[1]+50, box[0]+30, box[1]+38, box[0]+40, box[1]+44, box[0]+48, box[1]+28), fill=deep_blue, width=3)
            d.polygon([(box[0]+44, box[1]+28), (box[0]+50, box[1]+30), (box[0]+48, box[1]+24)], fill=deep_blue)
        else:
            d.ellipse((box[0]+18, box[1]+24, box[0]+46, box[1]+52), outline=deep_blue, width=3)
            d.line((box[0]+32, box[1]+24, box[0]+32, box[1]+52), fill=deep_blue, width=2)
            d.line((box[0]+18, box[1]+38, box[0]+46, box[1]+38), fill=deep_blue, width=2)
        draw_fitted_wrapped(d, (box[0]+62, box[1]+16), headline, 12, True, gray, box[2]-box[0]-72, 16, line_gap=1, min_size=9, max_lines=1)
        draw_fitted_wrapped(d, (box[0]+62, box[1]+34), body, 14, True, deep_blue, box[2]-box[0]-72, 34, line_gap=2, min_size=10, max_lines=2)
    draw_fitted_wrapped(d, (694, top_y+348), market_source_text(market_info), 10, False, gray, 492, 16, line_gap=1, min_size=8, max_lines=1)

    sec_y = 694
    left_sec = (24, sec_y, 550, 1100)
    right_sec = (566, sec_y, 1216, 1100)
    for box in [left_sec, right_sec]:
        _draw_shadowed_card(d, box, radius=22, fill=(255,255,255), outline=line, width=2, shadow=True)
    d.rounded_rectangle((38, sec_y+16, 74, sec_y+52), radius=8, fill=mix(primary, (255,255,255), 0.90))
    draw_simple_feature_icon(d, (44, sec_y+22, 68, sec_y+46), 'monitor', color=deep_blue)
    d.text((86, sec_y+14), label(lang, 'overview'), font=load_font(26, True), fill=deep_blue)
    d.rounded_rectangle((580, sec_y+16, 616, sec_y+52), radius=18, fill=mix(primary, (255,255,255), 0.90))
    d.ellipse((588, sec_y+24, 608, sec_y+44), outline=deep_blue, width=2)
    d.text((628, sec_y+14), label(lang, 'diff'), font=load_font(26, True), fill=deep_blue)

    def draw_info_rows(x1, y1, w, rows, icon_kinds, highlight_color):
        row_h = 90
        for i, txt in enumerate(rows[:3]):
            yy = y1 + i*(row_h+8)
            d.rounded_rectangle((x1, yy, x1+w, yy+row_h), radius=14, fill=(255,255,255), outline=line, width=1)
            d.ellipse((x1+14, yy+18, x1+58, yy+62), fill=mix(highlight_color, (255,255,255), 0.06), outline=None)
            draw_simple_feature_icon(d, (x1+20, yy+24, x1+52, yy+56), icon_kinds[i % len(icon_kinds)], color=(255,255,255))
            draw_emphasis_bullets_fit(d, x1+76, yy+20, [txt], 16, black, highlight_color, w-92, 54, bullet='', line_gap=4, item_gap=2, min_size=11, max_lines_per_item=3)

    draw_info_rows(36, sec_y+76, 500, data.get('overview', [])[:3], ['monitor','database','network'], deep_blue)
    draw_info_rows(578, sec_y+76, 624, data.get('differentiation', [])[:3], ['shield','database','network'], deep_blue)

    comp_y = 1126
    rep_box = (24, comp_y, 514, 1438)
    comp_box = (550, comp_y, 1216, 1438)
    for box in [rep_box, comp_box]:
        _draw_shadowed_card(d, box, radius=22, fill=(255,255,255), outline=line, width=2, shadow=True)
    d.rounded_rectangle((36, comp_y+16, 74, comp_y+52), radius=8, fill=mix(primary, (255,255,255), 0.90))
    draw_simple_feature_icon(d, (42, comp_y+22, 68, comp_y+46), 'monitor', color=deep_blue)
    d.text((86, comp_y+14), label(lang, 'drawing'), font=load_font(24, True), fill=deep_blue)
    d.rounded_rectangle((564, comp_y+16, 602, comp_y+52), radius=8, fill=mix(primary, (255,255,255), 0.90))
    draw_simple_feature_icon(d, (570, comp_y+22, 596, comp_y+46), 'shield', color=deep_blue)
    d.text((614, comp_y+14), label(lang, 'competitiveness'), font=load_font(24, True), fill=deep_blue)
    rep = fit_image(rep_img, (420, 228), bg=(255,255,255), trim=True)
    im.paste(rep, (24 + (490-rep.width)//2, comp_y+92))

    lim_box = (562, comp_y+84, 1198, comp_y+190)
    adv_box = (562, comp_y+254, 1198, comp_y+380)
    _draw_shadowed_card(d, lim_box, radius=14, fill=blue_tint, outline=line, width=1, shadow=False)
    d.text((584, comp_y+98), label(lang, 'limitations'), font=load_font(20, True), fill=deep_blue)
    draw_emphasis_bullets_fit(d, 586, comp_y+130, data.get('limitations', [])[:2], 14, black, deep_blue, 590, 50, bullet='•', line_gap=3, item_gap=4, min_size=10, max_lines_per_item=2)
    draw_arrow_down(im, 880, comp_y+208, color=mix(primary, (255,255,255), 0.50))
    _draw_shadowed_card(d, adv_box, radius=14, fill=green_tint, outline=line, width=1, shadow=False)
    d.text((584, comp_y+268), label(lang, 'advantages'), font=load_font(20, True), fill=deep_blue)
    draw_emphasis_bullets_fit(d, 586, comp_y+300, data.get('technical_advantages', [])[:2], 14, black, deep_blue, 590, 60, bullet='•', line_gap=3, item_gap=4, min_size=10, max_lines_per_item=2)

    ip_y = 1464
    _draw_shadowed_card(d, (24, ip_y, 1216, 1634), radius=22, fill=(255,255,255), outline=line, width=2, shadow=True)
    d.rounded_rectangle((36, ip_y+16, 74, ip_y+52), radius=8, fill=mix(primary, (255,255,255), 0.90))
    draw_simple_feature_icon(d, (42, ip_y+22, 68, ip_y+46), 'shield', color=deep_blue)
    d.text((86, ip_y+14), label(lang, 'ip'), font=load_font(24, True), fill=deep_blue)
    ip = normalize_ip(data.get('ip', {}))
    table_x, table_y = 44, ip_y+72
    table_w, table_h = 1152, 110
    header_h = 38
    col1 = int(table_w*0.42)
    col2 = int(table_w*0.25)
    c1, c2 = table_x+col1, table_x+col1+col2
    d.rounded_rectangle((table_x, table_y, table_x+table_w, table_y+table_h), radius=10, fill=(255,255,255), outline=line, width=1)
    d.rounded_rectangle((table_x, table_y, table_x+table_w, table_y+header_h), radius=10, fill=deep_blue)
    d.rectangle((table_x, table_y+20, table_x+table_w, table_y+header_h), fill=deep_blue)
    for xx in [c1, c2]:
        d.line((xx, table_y, xx, table_y+table_h), fill=line, width=1)
    draw_centered_wrapped(d, (table_x+8, table_y+4, c1-8, table_y+header_h-2), label(lang, 'invention'), load_font(15, True), (255,255,255), max_lines=1)
    draw_centered_wrapped(d, (c1+8, table_y+4, c2-8, table_y+header_h-2), label(lang, 'app_no'), load_font(12, True), (255,255,255), max_lines=2, line_gap=1)
    draw_centered_wrapped(d, (c2+8, table_y+4, table_x+table_w-8, table_y+header_h-2), label(lang, 'app_date'), load_font(12, True), (255,255,255), max_lines=2, line_gap=1)
    draw_centered_wrapped(d, (table_x+12, table_y+header_h+10, c1-12, table_y+table_h-8), ip['title'] or data.get('original_title', ''), load_font(14, False), black, max_lines=2, line_gap=2)
    num_text = f"{ip['application_number']}\n({ip['registration_number']})" if ip['registration_number'] else ip['application_number']
    date_text = f"{ip['application_date']}\n({ip['registration_date']})" if ip['registration_date'] else ip['application_date']
    draw_centered_wrapped(d, (c1+10, table_y+header_h+8, c2-10, table_y+table_h-8), num_text, load_font(15, False), black, max_lines=2, line_gap=2)
    draw_centered_wrapped(d, (c2+10, table_y+header_h+8, table_x+table_w-10, table_y+table_h-8), date_text, load_font(15, False), black, max_lines=2, line_gap=2)

    footer_y = 1650
    footer = Image.new('RGB', (1192, 78), (255,255,255))
    draw_linear_gradient(footer, (0,0,1192,78), deep_blue, mix(accent, (25,25,40), 0.25), horizontal=True)
    im.paste(footer, (24, footer_y))
    d.rounded_rectangle((24, footer_y, 1216, footer_y+78), radius=22, outline=line, width=1)
    d.ellipse((42, footer_y+18, 92, footer_y+68), fill=(255,255,255))
    d.text((58, footer_y+29), '☎', font=load_font(24, True), fill=deep_blue)
    d.text((106, footer_y+22), label(lang, 'contact'), font=load_font(24, True), fill=(255,255,255))
    draw_fitted_wrapped(d, (230, footer_y+28), contact, 18, False, (255,255,255), 940, 22, line_gap=2, min_size=12, max_lines=1)

    return im




def resolve_premium_reference_path() -> str:
    candidates = []
    if "__file__" in globals():
        candidates.append(os.path.join(os.path.dirname(__file__), PREMIUM_REFERENCE_IMAGE_PATH))
    candidates.append(os.path.join(os.getcwd(), PREMIUM_REFERENCE_IMAGE_PATH))
    candidates.append(PREMIUM_REFERENCE_IMAGE_PATH)
    for p in candidates:
        if p and os.path.exists(p):
            return p
    raise FileNotFoundError(f"Premium reference image not found: {PREMIUM_REFERENCE_IMAGE_PATH}")


def build_premium_infographic_ai_prompt(data: Dict[str, Any], contact: str, university_logo: Image.Image | None = None) -> str:
    lang = get_lang_code(data.get("language", "ko"))
    apps = data.get("applications", [])[:3]
    market = normalize_market_info(data.get("market_info", {}))
    app_lines = []
    for i, app in enumerate(apps, 1):
        app_lines.append(f"- 카드{i} 제목: {app.get('name','')}\n  카드{i} 설명: {app.get('description','')}")
    apps_text = "\n".join(app_lines)
    overview_text = "\n".join([f"- {x}" for x in data.get("overview", [])[:3]])
    diff_text = "\n".join([f"- {x}" for x in data.get("differentiation", [])[:3]])
    lim_text = "\n".join([f"- {x}" for x in data.get("limitations", [])[:2]])
    adv_text = "\n".join([f"- {x}" for x in data.get("technical_advantages", [])[:2]])
    ip = normalize_ip(data.get("ip", {}))
    title = data.get('marketing_title', '')
    subtitle = data.get('subtitle', '')
    university = data.get('university_display') or data.get('university', '')
    department = data.get('department_display') or data.get('department', '')
    professor = data.get('professor', '')
    prof_suffix = label(lang, 'prof_suffix')
    header_kicker = f"PIUM Tech Offer x {university} | {department} | {professor} {prof_suffix}"
    market_title = market.get('display_title') or market.get('market_name') or label(lang, 'market_default')
    years, values, cagr, unit = build_market_series(market)
    end_year = years[-1] if years else _safe_int(market.get('end_year'), 0)
    end_value = format_market_value(values[-1], unit) if values else str(market.get('end_value') or '-')
    theme = extract_logo_theme(university_logo)
    brand_primary = theme.get('primary', (0,55,135))
    brand_secondary = theme.get('secondary', mix(brand_primary, (30,115,210), 0.25))
    brand_accent = theme.get('accent', mix(brand_primary, (0,175,190), 0.35))
    brand_primary_hex = rgb_to_hex(brand_primary)
    brand_secondary_hex = rgb_to_hex(brand_secondary)
    brand_accent_hex = rgb_to_hex(brand_accent)
    prompt = f"""
You are given multiple reference images:
- Image A: premium_infographic_reference.png. Use this as the MAIN design/style reference.
- Image B: the current generated SMK page. Use this for the actual content, existing logos/QR block placement, representative drawing context, and information architecture.
- Image C: the original university logo asset. Use this exact logo visually; do not redraw it into a different logo.
- Image D: the original PIUM + QR reference card/block. Use this exact PIUM/QR block visually; do not create a second fake QR and do not duplicate the QR elsewhere.
- Image E: the original representative drawing. Use this drawing visually as the representative drawing panel; do not replace it with a different diagram.

Create ONE premium Korean technology infographic image based on Image A's layout quality and Image B's actual content.
This is NOT a blank redesign. It is a premium restyling of the current SMK using the supplied assets.

IMPORTANT ASSET + COLOR RULES:
- The university logo's dominant color is the authoritative brand color for the whole premium infographic.
- Do not let the generic blue reference override the university color identity.
- Preserve the university logo faithfully and keep it clean, sharp, and undistorted in the top-left area.
- Preserve the PIUM+QR block faithfully in the top-right area. Do not redraw it into a fake logo or fake QR code.
- Do not create extra QR codes. There should be only one PIUM/QR block in the top-right area.
- Preserve the representative drawing faithfully in the representative drawing section. Do not replace it with a different diagram.
- The final premium image itself must already contain the logo, PIUM+QR block, and representative drawing correctly. Do not rely on any later overlay or post-compositing step.

Design target:
- Visually close to Image A in layout quality and premium brochure polish, BUT do NOT copy Image A's blue color palette blindly.
- The main color palette MUST follow the university logo / Image C brand color.
- Use this brand palette as the dominant UI color system:
  - primary: {brand_primary_hex}
  - secondary: {brand_secondary_hex}
  - accent: {brand_accent_hex}
- All major section headers, title color, chart bars, footer, KPI icons, card accent strokes, and application-card accent colors should harmonize with this university brand palette.
- If the university logo is red, make the premium infographic red/pink/warm-accented like the original university identity; if blue, use blue; if green, use green. Do not force blue for every university.
- Stronger and more premium than the basic SMK.
- No awkward overlaps, no duplicated logos/QR, no cropped section panels.
- Keep all section order: header, applications/products, market status, technology overview, key differentiators, representative drawing, technical competitiveness, IP status, contact footer.
- Application/product cards should be full-color and visually rich.
- Market status should use a year-by-year bar chart with KPI cards.
- Overview and differentiators should use clean row cards with key phrases emphasized.
- Technical competitiveness should show current limitations and technical advantages clearly.
- IP status should be a neat table.
- Contact should be a premium footer bar.

Text/content to use:
Header line: {header_kicker}
Main title: {title}
Subtitle: {subtitle}

Applications / Products:
{apps_text}

Market section title: {market_title}
Representative market data summary:
- CAGR: {cagr:.1f}%
- End year: {end_year}
- End value: {end_value}
- Unit: {unit}
- Source: {market_source_text(market)}

Technology overview:
{overview_text}

Key differentiators:
{diff_text}

Current limitations:
{lim_text}

Technical advantages:
{adv_text}

IP section:
- Invention title: {ip.get('title','') or data.get('original_title','')}
- Application/Registration number: {ip.get('application_number','')} / {ip.get('registration_number','')}
- Application/Registration date: {ip.get('application_date','')} / {ip.get('registration_date','')}

Contact footer:
{contact}

Readability rules:
- Keep Korean text readable and professionally typeset.
- Avoid fake random Korean text.
- Avoid text truncation where possible.
- Preserve numbers, emails, phone numbers, and patent numbers as accurately as possible.
- Use concise Korean where needed for layout, but do not change the meaning.
"""
    return prompt


def _extract_generated_image_bytes(image_result: Any) -> bytes:
    # For images.generate / images.edit style responses
    data = getattr(image_result, 'data', None)
    if data:
        item = data[0]
        b64 = getattr(item, 'b64_json', None) or (item.get('b64_json') if isinstance(item, dict) else None)
        if b64:
            return base64.b64decode(b64)
    raise ValueError('No generated image bytes found in response')




def _save_temp_png(img: Image.Image, prefix: str) -> str:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.png', prefix=prefix)
    tmp.close()
    img.convert('RGBA' if img.mode == 'RGBA' else 'RGB').save(tmp.name, 'PNG')
    return tmp.name


def make_pium_qr_reference_card(pium_logo: Image.Image | None, qr_img: Image.Image | None, piumlink_logo: Image.Image | None = None) -> Image.Image:
    """프리미엄 AI 생성에 reference로 넘길 PIUM+QR 카드 블록.
    후합성용이 아니라, 모델이 원본 블록을 보고 자연스럽게 활용하도록 제공하는 자산이다.
    """
    card = Image.new('RGB', (360, 480), (255, 255, 255))
    d = ImageDraw.Draw(card)
    d.rounded_rectangle((8, 8, 352, 472), radius=34, fill=(255,255,255), outline=(190, 210, 232), width=4)
    if pium_logo is not None:
        logo = make_transparent_logo_canvas(pium_logo, size=(250, 88), padding=0)
        card.paste(logo, (55, 34), logo if logo.mode == 'RGBA' else None)
    # QR source priority: extracted QR > PIUMLINK bundled image
    qr_source = qr_img or piumlink_logo
    if qr_source is not None:
        qr = fit_image(qr_source.convert('RGB'), (250, 250), bg=(255,255,255), trim=True)
        card.paste(qr, (55, 145))
    return card



def make_brand_palette_reference_card(university_logo: Image.Image | None = None) -> Image.Image:
    """gpt-image-2에 대학 로고 기반 팔레트를 명확하게 전달하기 위한 색상 참조 카드."""
    theme = extract_logo_theme(university_logo)
    primary = theme.get("primary", (0,55,135))
    secondary = theme.get("secondary", mix(primary, (30,115,210), 0.25))
    accent = theme.get("accent", mix(primary, (0,175,190), 0.35))
    pale = theme.get("pale", mix(primary, (255,255,255), 0.88))
    card = Image.new('RGB', (640, 360), (255,255,255))
    d = ImageDraw.Draw(card)
    d.rounded_rectangle((12, 12, 628, 348), radius=28, fill=(255,255,255), outline=(210,220,232), width=3)
    d.text((34, 30), "UNIVERSITY BRAND PALETTE", font=load_font(30, True), fill=ensure_dark(primary))
    d.text((34, 70), "Use these colors for the premium infographic UI", font=load_font(20, False), fill=(90,98,112))
    colors = [("PRIMARY", primary), ("SECONDARY", secondary), ("ACCENT", accent), ("SOFT BG", pale)]
    x = 38
    for label_text, color in colors:
        d.rounded_rectangle((x, 120, x+128, 238), radius=18, fill=color, outline=(230,235,242), width=2)
        d.text((x, 254), label_text, font=load_font(16, True), fill=(40,45,55))
        d.text((x, 278), rgb_to_hex(color), font=load_font(16, False), fill=(80,88,100))
        x += 148
    if university_logo is not None:
        logo = make_transparent_logo_canvas(university_logo, size=(92,92), padding=0)
        card.paste(logo, (508, 30), logo if logo.mode == 'RGBA' else None)
    return card



def _resample_qr(img: Image.Image, size: tuple[int,int]) -> Image.Image:
    src = img.convert('RGB')
    return ImageOps.contain(src, size, Image.NEAREST)


def _resample_logo(img: Image.Image, size: tuple[int,int]) -> Image.Image:
    if img.mode != 'RGBA':
        src = img.convert('RGBA')
    else:
        src = img.copy()
    fitted = make_transparent_logo_canvas(src, size=size, padding=0)
    # very light sharpening improves small text/logo edges after downscaling
    return fitted.filter(ImageFilter.UnsharpMask(radius=0.8, percent=110, threshold=2))


def composite_brand_assets_on_final_image(
    final_img: Image.Image,
    university_logo: Image.Image | None = None,
    pium_logo: Image.Image | None = None,
    qr_img: Image.Image | None = None,
    piumlink_logo: Image.Image | None = None,
) -> Image.Image:
    """최종 결과물 위에 원본 대학 로고 / PIUM 로고 / QR 을 마지막 단계에서 다시 정합 합성.
    목적: 로고/QR의 선명도 보존, 스캔 안정성 확보.
    레이아웃은 현재 compose_infographic_brief의 기준 좌표를 비율로 환산해 정합한다.
    """
    im = final_img.convert('RGBA').copy()
    W, H = im.size
    sx = W / 1240.0
    sy = H / 1754.0
    d = ImageDraw.Draw(im)
    line = (196, 212, 228, 255)
    # top-left university logo area
    if university_logo is not None:
        lx, ly = int(round(28*sx)), int(round(24*sy))
        lw, lh = int(round(118*sx)), int(round(118*sy))
        # clear only the logo content region, preserving the surrounding layout
        d.rounded_rectangle((lx-6, ly-6, lx+lw+6, ly+lh+6), radius=max(12, int(18*sx)), fill=(255,255,255,0))
        logo = _resample_logo(university_logo, (lw, lh))
        im.paste(logo, (lx, ly), logo if logo.mode == 'RGBA' else None)

    # top-right PIUM + QR card area
    qx1, qy1 = int(round(1016*sx)), int(round(20*sy))
    qx2, qy2 = int(round(1210*sx)), int(round(264*sy))
    # redraw card cleanly for crisp contents while respecting layout position
    d.rounded_rectangle((qx1, qy1, qx2, qy2), radius=max(18, int(24*sx)), fill=(255,255,255,255), outline=line, width=max(2, int(2*sx)))
    if pium_logo is not None:
        pl = _resample_logo(pium_logo, (int(round(124*sx)), int(round(48*sy))))
        px, py = int(round(1052*sx)), int(round(36*sy))
        im.paste(pl, (px, py), pl if pl.mode == 'RGBA' else None)
    qr_source = qr_img or piumlink_logo
    if qr_source is not None:
        # prefer extracted pure QR image; if the fallback is a broader PIUMLINK image,
        # fit it tightly but still use nearest-neighbor for crisp square edges.
        qrs = _resample_qr(qr_source, (int(round(134*sx)), int(round(134*sy))))
        # center if the aspect result is smaller than target
        bg = Image.new('RGB', (int(round(134*sx)), int(round(134*sy))), 'white')
        bg.paste(qrs, ((bg.width-qrs.width)//2, (bg.height-qrs.height)//2))
        qx, qy = int(round(1047*sx)), int(round(92*sy))
        im.paste(bg, (qx, qy))
    return im.convert('RGB')

def overlay_preserved_assets_on_premium_infographic(
    base_img: Image.Image,
    university_logo: Image.Image | None = None,
    pium_logo: Image.Image | None = None,
    qr_img: Image.Image | None = None,
    rep_img: Image.Image | None = None,
    piumlink_logo: Image.Image | None = None,
) -> Image.Image:
    im = composite_brand_assets_on_final_image(
        base_img,
        university_logo=university_logo,
        pium_logo=pium_logo,
        qr_img=qr_img,
        piumlink_logo=piumlink_logo,
    )
    d = ImageDraw.Draw(im)
    line = (200, 214, 230)
    if rep_img is not None:
        sx = im.width / 1240.0
        sy = im.height / 1754.0
        rep_box = (int(round(34*sx)), int(round(975*sy)), int(round(424*sx)), int(round(1282*sy)))
        d.rounded_rectangle(rep_box, radius=max(18, int(20*sx)), fill=(255,255,255), outline=line, width=max(2, int(2*sx)))
        rep = fit_image(rep_img.convert('RGB'), (rep_box[2]-rep_box[0]-30, rep_box[3]-rep_box[1]-32), bg=(255,255,255), trim=True)
        im.paste(rep, (rep_box[0] + (rep_box[2]-rep_box[0]-rep.width)//2, rep_box[1] + (rep_box[3]-rep_box[1]-rep.height)//2))
    return im




def _draw_exact_header_overlay(img: Image.Image, data: Dict[str, Any], university_logo: Image.Image | None = None) -> Image.Image:
    im = img.convert('RGBA').copy()
    d = ImageDraw.Draw(im)
    W, H = im.size
    sx, sy = W/1240.0, H/1754.0
    theme = extract_logo_theme(university_logo)
    primary = ensure_dark(theme.get('primary', (0,55,135)))
    black = (30,36,46,255)
    gray = (88,96,108,255)
    # clean exact text area while preserving overall background outside
    x1, y1 = int(176*sx), int(30*sy)
    x2, y2 = int(990*sx), int(238*sy)
    d.rounded_rectangle((x1, y1, x2, y2), radius=max(16,int(22*sx)), fill=(255,255,255,238))
    lang = get_lang_code(data.get('language','ko'))
    kicker = f"PIUM Tech Offer x {data.get('university_display') or data.get('university','')} | {data.get('department_display') or data.get('department','')} | {data.get('professor','')} {label(lang,'prof_suffix')}"
    draw_fitted_wrapped(d, (x1+int(18*sx), y1+int(16*sy)), kicker, max(12,int(16*sx)), False, gray[:3], x2-x1-int(36*sx), int(24*sy), line_gap=2, min_size=max(9,int(10*sx)), max_lines=1)
    draw_fitted_wrapped(d, (x1+int(18*sx), y1+int(48*sy)), str(data.get('marketing_title','')), max(24,int(38*sx)), True, primary, x2-x1-int(36*sx), int(92*sy), line_gap=max(2,int(4*sy)), min_size=max(18,int(20*sx)), max_lines=3)
    subtitle = str(data.get('subtitle') or '').strip()
    if subtitle:
        draw_fitted_wrapped(d, (x1+int(18*sx), y1+int(150*sy)), subtitle, max(12,int(18*sx)), False, black[:3], x2-x1-int(36*sx), int(38*sy), line_gap=2, min_size=max(9,int(10*sx)), max_lines=2)
    return im.convert('RGB')


def _draw_exact_ip_table_overlay(img: Image.Image, data: Dict[str, Any], university_logo: Image.Image | None = None) -> Image.Image:
    im = img.convert('RGBA').copy()
    d = ImageDraw.Draw(im)
    W,H=im.size
    sx, sy = W/1240.0, H/1754.0
    theme=extract_logo_theme(university_logo)
    primary=ensure_dark(theme.get('primary',(0,55,135)))
    line=mix(primary,(210,220,232),0.82)
    black=(30,36,46)
    lang=get_lang_code(data.get('language','ko'))
    ip=normalize_ip(data.get('ip',{}))
    # overwrite lower IP section with exact table
    box=(int(24*sx), int(1460*sy), int(1216*sx), int(1638*sy))
    d.rounded_rectangle(box, radius=max(16,int(22*sx)), fill=(255,255,255,246), outline=line, width=max(1,int(2*sx)))
    d.text((box[0]+int(30*sx), box[1]+int(16*sy)), label(lang,'ip'), font=load_font(max(16,int(24*sx)), True), fill=primary)
    table_x, table_y = box[0]+int(20*sx), box[1]+int(64*sy)
    table_w, table_h = box[2]-box[0]-int(40*sx), int(94*sy)
    header_h=int(32*sy)
    col1=int(table_w*0.44); col2=int(table_w*0.27)
    c1=table_x+col1; c2=c1+col2
    d.rounded_rectangle((table_x,table_y,table_x+table_w,table_y+table_h), radius=max(8,int(10*sx)), fill=(255,255,255,255), outline=line, width=max(1,int(1*sx)))
    d.rectangle((table_x,table_y,table_x+table_w,table_y+header_h), fill=primary)
    for xx in (c1,c2):
        d.line((xx, table_y, xx, table_y+table_h), fill=line, width=max(1,int(1*sx)))
    white=(255,255,255)
    draw_centered_wrapped(d,(table_x+5,table_y+3,c1-5,table_y+header_h-2),label(lang,'invention'),load_font(max(10,int(14*sx)),True),white,max_lines=1)
    draw_centered_wrapped(d,(c1+5,table_y+2,c2-5,table_y+header_h-2),label(lang,'app_no'),load_font(max(8,int(11*sx)),True),white,max_lines=2,line_gap=1)
    draw_centered_wrapped(d,(c2+5,table_y+2,table_x+table_w-5,table_y+header_h-2),label(lang,'app_date'),load_font(max(8,int(11*sx)),True),white,max_lines=2,line_gap=1)
    title=ip.get('title') or data.get('original_title','')
    num = f"{ip.get('application_number','')}\n({ip.get('registration_number','')})" if ip.get('registration_number') else ip.get('application_number','')
    date = f"{ip.get('application_date','')}\n({ip.get('registration_date','')})" if ip.get('registration_date') else ip.get('application_date','')
    draw_centered_wrapped(d,(table_x+8,table_y+header_h+5,c1-8,table_y+table_h-5),title,load_font(max(9,int(13*sx)),False),black,max_lines=2,line_gap=2)
    draw_centered_wrapped(d,(c1+8,table_y+header_h+5,c2-8,table_y+table_h-5),num,load_font(max(9,int(13*sx)),False),black,max_lines=2,line_gap=2)
    draw_centered_wrapped(d,(c2+8,table_y+header_h+5,table_x+table_w-8,table_y+table_h-5),date,load_font(max(9,int(13*sx)),False),black,max_lines=2,line_gap=2)
    return im.convert('RGB')


def _draw_exact_contact_overlay(img: Image.Image, contact: str, data: Dict[str, Any], university_logo: Image.Image | None = None) -> Image.Image:
    im=img.convert('RGBA').copy()
    d=ImageDraw.Draw(im)
    W,H=im.size; sx,sy=W/1240.0,H/1754.0
    theme=extract_logo_theme(university_logo)
    primary=ensure_dark(theme.get('primary',(0,55,135)))
    accent=theme.get('accent', mix(primary,(0,175,190),0.35))
    footer=(int(24*sx), int(1650*sy), int(1216*sx), int(1730*sy))
    # gradient footer
    w=max(1, footer[2]-footer[0]); h=max(1, footer[3]-footer[1])
    grad=Image.new('RGB',(w,h),primary)
    draw_linear_gradient(grad,(0,0,w,h),primary,mix(accent,(25,25,40),0.20),horizontal=True)
    im.paste(grad.convert('RGBA'),(footer[0],footer[1]))
    d.rounded_rectangle(footer, radius=max(16,int(22*sx)), outline=(210,220,232,255), width=max(1,int(1*sx)))
    lang=get_lang_code(data.get('language','ko'))
    d.text((footer[0]+int(34*sx), footer[1]+int(22*sy)), label(lang,'contact'), font=load_font(max(15,int(24*sx)),True), fill=(255,255,255,255))
    draw_fitted_wrapped(d,(footer[0]+int(160*sx), footer[1]+int(27*sy)), contact, max(11,int(17*sx)), False, (255,255,255), footer[2]-footer[0]-int(200*sx), int(26*sy), line_gap=2, min_size=max(8,int(10*sx)), max_lines=1)
    return im.convert('RGB')



PREMIUM_SECTION_BBOX = {
    "header": (24, 20, 1216, 264),
    "apps_market": (24, 280, 1216, 670),
    "overview_diff": (24, 694, 1216, 1100),
    "rep_comp": (24, 1126, 1216, 1438),
    "ip": (24, 1460, 1216, 1638),
    "footer": (24, 1650, 1216, 1730),
}

SECTION_FALLBACK_ORDER = ["apps_market", "overview_diff", "rep_comp"]


def _scale_box(box: tuple[int,int,int,int], img_size: tuple[int,int]) -> tuple[int,int,int,int]:
    W, H = img_size
    sx, sy = W / 1240.0, H / 1754.0
    x1, y1, x2, y2 = box
    return (int(round(x1 * sx)), int(round(y1 * sy)), int(round(x2 * sx)), int(round(y2 * sy)))


def _crop_box(img: Image.Image, box: tuple[int,int,int,int]) -> Image.Image:
    return img.crop(box)


def _gray_stats(img: Image.Image) -> tuple[float, float, float]:
    arr = np.asarray(img.convert('L'), dtype=np.float32)
    mean = float(arr.mean()) if arr.size else 0.0
    std = float(arr.std()) if arr.size else 0.0
    dark_ratio = float((arr < 180).mean()) if arr.size else 0.0
    return mean, std, dark_ratio


def _edge_density(img: Image.Image) -> float:
    arr = np.asarray(img.convert('L'), dtype=np.float32)
    if arr.ndim != 2 or arr.shape[0] < 2 or arr.shape[1] < 2:
        return 0.0
    gx = np.abs(np.diff(arr, axis=1))
    gy = np.abs(np.diff(arr, axis=0))
    edge = (gx > 18).mean() * 0.5 + (gy > 18).mean() * 0.5
    return float(edge)


def _text_like_score(img: Image.Image) -> float:
    mean, std, dark_ratio = _gray_stats(img)
    edge = _edge_density(img)
    # heuristic: higher std + enough dark pixels + enough edges implies readable content
    return std * 0.04 + dark_ratio * 2.2 + edge * 6.0


def _should_replace_section(generated_crop: Image.Image, reference_crop: Image.Image, section_name: str) -> bool:
    g_score = _text_like_score(generated_crop)
    r_score = _text_like_score(reference_crop)
    g_mean, g_std, g_dark = _gray_stats(generated_crop)
    # blank / washed / extremely low-detail detection
    if g_std < 10 or g_dark < 0.02:
        return True
    # if generated section is much weaker than stable section, fallback that section
    if g_score < r_score * 0.55:
        return True
    # additional stricter rule for text-heavy sections
    if section_name in ("overview_diff", "rep_comp") and g_score < r_score * 0.68:
        return True
    return False


def _paste_section(base_img: Image.Image, section_img: Image.Image, box: tuple[int,int,int,int]) -> Image.Image:
    im = base_img.copy()
    target = (box[2] - box[0], box[3] - box[1])
    if section_img.size != target:
        section_img = section_img.resize(target, Image.LANCZOS)
    im.paste(section_img, (box[0], box[1]))
    return im


def build_stable_premium_reference_image(
    data: Dict[str, Any],
    rep_img: Image.Image,
    app_imgs: List[Image.Image],
    contact: str,
    university_logo: Image.Image | None = None,
    pium_logo: Image.Image | None = None,
    qr_img: Image.Image | None = None,
    piumlink_logo: Image.Image | None = None,
) -> Image.Image:
    return make_high_quality_infographic_image(
        data,
        rep_img,
        app_imgs,
        contact,
        university_logo=university_logo,
        pium_logo=pium_logo,
        qr_img=qr_img,
        piumlink_logo=piumlink_logo,
    )


def apply_section_validation_and_corrections(
    generated: Image.Image,
    stable_reference: Image.Image,
) -> Image.Image:
    """섹션 단위 자동 검증 + 보정.
    현재는 텍스트/레이아웃 오류가 잦은 본문 섹션을 대상으로,
    생성 결과가 지나치게 빈약하거나 읽기 어려운 경우 안정형 섹션으로 교체한다.
    """
    img = generated.convert('RGB')
    ref = stable_reference.convert('RGB')
    for name in SECTION_FALLBACK_ORDER:
        box = _scale_box(PREMIUM_SECTION_BBOX[name], img.size)
        g_crop = _crop_box(img, box)
        r_crop = _crop_box(ref, box)
        if _should_replace_section(g_crop, r_crop, name):
            img = _paste_section(img, r_crop, box)
    return img



def _detect_rep_drawing_bbox(img: Image.Image) -> tuple[int,int,int,int] | None:
    """생성 이미지 안의 대표도면 패널 위치를 대략 탐지한다.
    고정 섹션 전체를 덮지 않고, 실제 도면이 있는 하단 좌측 영역의 어두운 선분 분포만 이용한다.
    """
    W, H = img.size
    # lower-left region where representative drawing is normally located
    rx1, ry1 = int(W * 0.03), int(H * 0.58)
    rx2, ry2 = int(W * 0.47), int(H * 0.82)
    roi = img.crop((rx1, ry1, rx2, ry2)).convert('L')
    arr = np.asarray(roi, dtype=np.uint8)
    # ignore very top part of ROI to avoid section title
    top_ignore = int(arr.shape[0] * 0.12)
    work = arr[top_ignore:, :]
    # black/gray technical drawing lines
    mask = work < 130
    # remove sparse noise by row/column support
    if mask.mean() < 0.003:
        return None
    ys, xs = np.where(mask)
    if len(xs) < 80 or len(ys) < 80:
        return None
    x1, x2 = int(xs.min()), int(xs.max())
    y1, y2 = int(ys.min()) + top_ignore, int(ys.max()) + top_ignore
    bw, bh = x2 - x1, y2 - y1
    if bw < W * 0.12 or bh < H * 0.06:
        return None
    # clamp unrealistically large detections
    if bw > W * 0.44 or bh > H * 0.28:
        return None
    pad_x = int(W * 0.025)
    pad_y = int(H * 0.025)
    bx1 = max(rx1, rx1 + x1 - pad_x)
    by1 = max(ry1, ry1 + y1 - pad_y)
    bx2 = min(rx2, rx1 + x2 + pad_x)
    by2 = min(ry2, ry1 + y2 + pad_y)
    # ensure a reasonably rectangular drawing box
    if bx2 - bx1 < W * 0.18 or by2 - by1 < H * 0.10:
        return None
    return (bx1, by1, bx2, by2)


def _default_rep_drawing_bbox(img: Image.Image) -> tuple[int,int,int,int]:
    W, H = img.size
    sx, sy = W / 1240.0, H / 1754.0
    return (int(round(62*sx)), int(round(1010*sy)), int(round(458*sx)), int(round(1295*sy)))


def overlay_rep_drawing_micro_patch(img: Image.Image, rep_img: Image.Image | None) -> Image.Image:
    """기존 생성 결과의 대표도면 위치를 최대한 유지하면서, 도면 그림 부분만 원본으로 교체한다."""
    if rep_img is None:
        return img.convert('RGB')
    im = img.convert('RGB').copy()
    bbox = _detect_rep_drawing_bbox(im) or _default_rep_drawing_bbox(im)
    x1, y1, x2, y2 = bbox
    d = ImageDraw.Draw(im)
    # only clear the detected drawing area, not the whole section/card
    d.rounded_rectangle((x1, y1, x2, y2), radius=max(8, int((x2-x1)*0.035)), fill=(255,255,255), outline=(205,216,228), width=1)
    margin_x = max(8, int((x2-x1)*0.05))
    margin_y = max(8, int((y2-y1)*0.06))
    rep = fit_image(rep_img.convert('RGB'), (max(1, x2-x1-2*margin_x), max(1, y2-y1-2*margin_y)), bg=(255,255,255), trim=True)
    im.paste(rep, (x1 + (x2-x1-rep.width)//2, y1 + (y2-y1-rep.height)//2))
    return im

def apply_mandatory_premium_corrections(
    generated: Image.Image,
    data: Dict[str, Any],
    rep_img: Image.Image,
    contact: str,
    university_logo: Image.Image | None = None,
    pium_logo: Image.Image | None = None,
    qr_img: Image.Image | None = None,
    piumlink_logo: Image.Image | None = None,
) -> Image.Image:
    """실제 사용용 1차 자동 보정 파이프라인.
    AI 생성 결과의 미감은 최대한 유지하되, 치명 오류가 자주 나는 영역은 JSON/원본 자산으로 확정 보정한다.
    """
    img = generated.convert('RGB')
    img = composite_brand_assets_on_final_image(img, university_logo=university_logo, pium_logo=pium_logo, qr_img=qr_img, piumlink_logo=piumlink_logo)
    img = _draw_exact_header_overlay(img, data, university_logo=university_logo)
    # 대표도면은 생성 결과의 도면 패널을 탐지한 뒤 해당 그림 영역만 micro patch
    img = overlay_rep_drawing_micro_patch(img, rep_img)
    img = _draw_exact_ip_table_overlay(img, data, university_logo=university_logo)
    img = _draw_exact_contact_overlay(img, contact, data, university_logo=university_logo)
    return img

def generate_reference_based_premium_infographic(
    data: Dict[str, Any],
    rep_img: Image.Image,
    app_imgs: List[Image.Image],
    contact: str,
    university_logo: Image.Image | None = None,
    pium_logo: Image.Image | None = None,
    qr_img: Image.Image | None = None,
    piumlink_logo: Image.Image | None = None,
    base_smk_img: Image.Image | None = None,
) -> Image.Image:
    """레퍼런스 기반 프리미엄 인포그래픽 생성.
    후합성으로 강제로 덮지 않고, 프리미엄 레퍼런스 + 현재 SMK + 원본 자산들을
    모두 입력 이미지로 제공한 뒤 gpt-image-2가 자연스럽게 활용하도록 한다.
    실패 시 규칙 기반 렌더러로 fallback한다.
    """
    temp_paths: List[str] = []
    try:
        client = get_client()
        ref_path = resolve_premium_reference_path()
        prompt = build_premium_infographic_ai_prompt(data, contact, university_logo=university_logo)

        # Reference images: style reference + current SMK + exact asset references.
        temp_paths.append(ref_path)
        if base_smk_img is not None:
            temp_paths.append(_save_temp_png(base_smk_img, 'premium_base_smk_'))
        if university_logo is not None:
            temp_paths.append(_save_temp_png(university_logo, 'premium_univ_logo_'))
        # Explicit brand palette reference: prevents the reference image's blue palette from overriding university-specific colors.
        temp_paths.append(_save_temp_png(make_brand_palette_reference_card(university_logo), 'premium_brand_palette_'))
        pium_qr_card = make_pium_qr_reference_card(pium_logo, qr_img, piumlink_logo)
        temp_paths.append(_save_temp_png(pium_qr_card, 'premium_pium_qr_'))
        if rep_img is not None:
            temp_paths.append(_save_temp_png(rep_img, 'premium_rep_drawing_'))

        with ExitStack() as stack:
            files = [stack.enter_context(open(p, 'rb')) for p in temp_paths]
            result = client.images.edit(
                model=PREMIUM_IMAGE_MODEL_FIXED,
                image=files,
                prompt=prompt,
                size='1024x1536',
            )
        raw = _extract_generated_image_bytes(result)
        generated = Image.open(BytesIO(raw)).convert('RGB')
        return generated
    except Exception:
        # fallback: preserve operational stability
        return make_high_quality_infographic_image(
            data,
            rep_img,
            app_imgs,
            contact,
            university_logo=university_logo,
            pium_logo=pium_logo,
            qr_img=qr_img,
            piumlink_logo=piumlink_logo,
        )
    finally:
        # remove temporary files except bundled reference path
        ref_abs = None
        try:
            ref_abs = os.path.abspath(resolve_premium_reference_path())
        except Exception:
            pass
        for p in temp_paths:
            try:
                if ref_abs and os.path.abspath(p) == ref_abs:
                    continue
                if os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass


def make_high_quality_infographic_image(




    data: Dict[str, Any],
    rep_img: Image.Image,
    app_imgs: List[Image.Image],
    contact: str,
    university_logo: Image.Image | None = None,
    pium_logo: Image.Image | None = None,
    qr_img: Image.Image | None = None,
    piumlink_logo: Image.Image | None = None,
) -> Image.Image:
    """규칙 기반의 고급 인포그래픽 렌더러.
    기본 SMK의 구조는 유지하되, 보다 고급스러운 카드/배경/타이포를 적용하고
    결과는 2배 해상도의 배포용 PNG/PDF로 출력한다.
    """
    premium = compose_infographic_brief(
        data,
        rep_img,
        app_imgs,
        contact,
        university_logo=university_logo,
        pium_logo=pium_logo,
        qr_img=qr_img,
        piumlink_logo=piumlink_logo,
    ).convert("RGB")

    scale = 2
    hq = premium.resize((premium.width * scale, premium.height * scale), Image.LANCZOS)
    hq = ImageEnhance.Contrast(hq).enhance(1.02)
    hq = ImageEnhance.Sharpness(hq).enhance(1.10)
    hq = hq.filter(ImageFilter.UnsharpMask(radius=1.0, percent=75, threshold=2))

    # 로고/QR은 최종 해상도 이미지에서 원본을 다시 정합 합성해 선명도와 형태를 보존
    hq = composite_brand_assets_on_final_image(
        hq,
        university_logo=university_logo,
        pium_logo=pium_logo,
        qr_img=qr_img,
        piumlink_logo=piumlink_logo,
    )

    return hq


def make_pdf_bytes_from_hq_image(img: Image.Image) -> bytes:
    """프리미엄 인포그래픽 PDF를 실제 A4 300dpi급 이미지형 PDF로 저장한다."""
    return _image_to_a4_pdf_bytes(img, dpi=300)


def make_png_bytes_from_image(img: Image.Image) -> bytes:
    bio = BytesIO()
    img.save(bio, format="PNG")
    return bio.getvalue()


def img_bytes(img: Image.Image) -> BytesIO:
    bio = BytesIO()
    img.save(bio, format="PNG")
    bio.seek(0)
    return bio

def px(x):
    return Inches(x / 124.0)  # 1240px -> 10in

def add_textbox(slide, x, y, w, h, text, size=14, bold=False, color=(30,30,30), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(1)
    tf.margin_right = Pt(1)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = str(text)
    run.font.name = "Noto Sans CJK KR"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box

def add_rect(slide, x, y, w, h, fill=(255,255,255), outline=(220,230,242), radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, px(x), px(y), px(w), px(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(*fill)
    shp.line.color.rgb = RGBColor(*outline)
    shp.line.width = Pt(1)
    return shp


def make_pptx_bytes(data: Dict[str, Any], rep_img: Image.Image, app_imgs: List[Image.Image], contact: str, university_logo: Image.Image | None = None, pium_logo: Image.Image | None = None, qr_img: Image.Image | None = None, piumlink_logo: Image.Image | None = None) -> bytes:
    """시장현황 포함 PPTX 생성."""
    pal = get_visual_palette(university_logo)
    uni_primary = pal["uni_primary"]
    uni_pale = pal["uni_pale"]
    uni_line = pal["uni_line"]
    primary = pal["pium_blue"]
    secondary = pal["tech_blue"]
    accent = pal["tech_cyan"]
    sky = pal["tech_pale"]
    sky2 = pal["tech_pale2"]
    line = pal["tech_line"]
    table_header = pal["table_header"]
    black = pal["black"]
    gray = pal["gray"]
    lang = get_lang_code(data.get("language", "ko"))

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(14.145)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 긴 기술명 대응: 이미지 렌더러와 동일하게 헤더 높이/좌표 확장
    add_rect(slide, 0, 0, 1240, 280, fill=uni_pale, outline=uni_pale)
    uni_logo_size = 124
    left_logo = make_transparent_logo_canvas(university_logo, size=(uni_logo_size, uni_logo_size), padding=0) if university_logo else make_university_logo_box(None, data.get('university_display') or data.get('university',''), size=(uni_logo_size, uni_logo_size), bg=uni_pale, primary=uni_primary)
    slide.shapes.add_picture(img_bytes(left_logo), px(28), px(24), width=px(uni_logo_size), height=px(uni_logo_size))

    # PIUM 로고/QR 보호 카드: 원본 로고 컬러는 유지하고, 대학색 헤더와는 시각적으로 분리
    right_card_x, right_card_y = 1240 - 178, 18
    right_card_w, right_card_h = 148, 220
    add_rect(slide, right_card_x + 5, right_card_y + 5, right_card_w, right_card_h, fill=mix(uni_line, (255,255,255), 0.55), outline=mix(uni_line, (255,255,255), 0.55), radius=True)
    add_rect(slide, right_card_x, right_card_y, right_card_w, right_card_h, fill=(255,255,255), outline=uni_line, radius=True)
    right_x = right_card_x + 10
    top_y = right_card_y + 14
    if pium_logo is not None:
        pium_canvas = make_transparent_logo_canvas(pium_logo, size=(128, 48), padding=0)
        slide.shapes.add_picture(img_bytes(pium_canvas), px(right_x), px(top_y), width=px(128), height=px(48))
    if piumlink_logo is not None:
        link_size = 112
        link = make_transparent_logo_canvas(piumlink_logo, size=(link_size,link_size), padding=0)
        slide.shapes.add_picture(img_bytes(link), px(right_card_x + (right_card_w-link_size)//2), px(top_y+60), width=px(link_size), height=px(link_size))

    header_x = 190
    header_w = right_card_x - header_x - 24
    prof_suffix = label(lang, "prof_suffix")
    add_textbox(slide, header_x, 42, header_w, 30, f"PIUM Tech Offer  x  {data.get('university_display') or data.get('university','')}  |  {data.get('department_display') or data.get('department','')}  |  {data.get('professor','')} {prof_suffix}", 13.5, False, uni_primary)
    add_textbox(slide, header_x, 78, header_w, 108, data.get("marketing_title", ""), 24, True, uni_primary)
    subtitle = (data.get("subtitle", "") or "").strip()
    if subtitle:
        sub_text = ensure_dark(mix(uni_primary, black, 0.35))
        add_textbox(slide, header_x, 204, header_w, 26, f"- {subtitle}", 11.2, False, sub_text, align=PP_ALIGN.LEFT)

    X = 28; CW = 1184

    # Applications + Market
    app_y, app_h = 294, 262
    top_gap = 30
    app_card_w = 552
    market_x = X + app_card_w + top_gap
    market_w = CW - app_card_w - top_gap

    add_rect(slide, X, app_y, app_card_w, app_h, fill=(255,255,255), outline=line, radius=True)
    add_rect(slide, market_x, app_y, market_w, app_h, fill=(255,255,255), outline=line, radius=True)
    add_textbox(slide, X+34, app_y+30, 240, 36, label(lang, "apps"), 15.5, True, primary)
    add_textbox(slide, market_x+34, app_y+30, 160, 36, label(lang, "market"), 15.5, True, primary)

    apps = data.get("applications", [])[:3]
    app_inner_x = X + 22
    app_inner_w = app_card_w - 44
    col_gap = 4
    col_w = (app_inner_w - col_gap*2) // 3
    for i in range(3):
        col_x = app_inner_x + i*(col_w + col_gap)
        if i < len(app_imgs):
            icon_w, icon_h = 112, 88
            slide.shapes.add_picture(img_bytes(fit_image(app_imgs[i], (icon_w, icon_h), trim=True)), px(col_x+(col_w-icon_w)//2), px(app_y+82), width=px(icon_w), height=px(icon_h))
        name = apps[i].get("name", "") if i < len(apps) else ""
        add_textbox(slide, col_x+6, app_y+176, col_w-12, 48, name, 11.5, True, black, align=PP_ALIGN.CENTER)

    market_info = normalize_market_info(data.get("market_info", {}))
    market_title = market_info.get("display_title") or market_info.get("market_name") or label(lang, "market_default")
    add_textbox(slide, market_x+34, app_y+74, market_w-68, 28, market_title, 11.2, True, black)

    graph_x, graph_y = market_x+34, app_y+104
    graph_w, graph_h = 178, 104
    add_rect(slide, graph_x, graph_y, graph_w, graph_h, fill=(255,255,255), outline=line, radius=True)
    chart, _market_visual_mode = get_market_visual(market_info, primary=primary, accent=accent)
    slide.shapes.add_picture(img_bytes(fit_image(chart, (graph_w-10, graph_h-10), trim=False)), px(graph_x+5), px(graph_y+5), width=px(graph_w-10), height=px(graph_h-10))

    desc_x = graph_x + graph_w + 16
    desc_y = app_y + 114
    desc_w = market_x + market_w - 28 - desc_x
    desc_h = 78
    add_rect(slide, desc_x, desc_y, desc_w, desc_h, fill=(255,255,255), outline=line, radius=True)
    add_textbox(slide, desc_x+12, desc_y+10, desc_w-24, desc_h-16, market_description_text(market_info), 7.9, False, gray)

    add_textbox(slide, graph_x, app_y+218, market_w-64, 16, market_source_text(market_info), 6.6, False, gray)

    # Overview / Difference
    y=582; gap=30; box_w=(CW-gap)//2; box_h=290
    for bx,title,items in [(X,label(lang, "overview"),data.get("overview",[])[:3]),(X+box_w+gap,label(lang, "diff"),data.get("differentiation",[])[:3])]:
        add_rect(slide, bx, y, box_w, box_h, fill=sky, outline=line, radius=True)
        add_textbox(slide, bx+36, y+34, 250, 34, title, 15, True, primary)
        add_textbox(slide, bx+42, y+88, box_w-84, box_h-112, "\n".join(["› "+str(v) for v in items]), 10.5, False, black)

    # Rep + competitiveness
    y=900; rep_w_box=380; comp_x=X+rep_w_box+22; comp_w=X+CW-comp_x; comp_h=322
    add_rect(slide, X, y, rep_w_box, comp_h, fill=(255,255,255), outline=line, radius=True)
    add_rect(slide, comp_x, y, comp_w, comp_h, fill=(255,255,255), outline=line, radius=True)
    add_textbox(slide, X+36, y+34, 180, 34, label(lang, "drawing"), 15, True, primary)
    slide.shapes.add_picture(img_bytes(fit_image(rep_img, (rep_w_box-74,214), trim=True)), px(X+37), px(y+82), width=px(rep_w_box-74), height=px(214))
    add_textbox(slide, comp_x+36, y+34, 220, 34, label(lang, "competitiveness"), 15, True, primary)
    inner_x=comp_x+44; inner_w=comp_w-70
    add_rect(slide, inner_x, y+76, inner_w, 102, fill=(255,255,255), outline=line, radius=True)
    add_textbox(slide, inner_x+20, y+90, 220, 22, label(lang, "limitations"), 10.5, True, gray)
    add_textbox(slide, inner_x+22, y+118, inner_w-44, 52, "\n".join(["• "+str(v) for v in data.get("limitations",[])[:2]]), 8.2, False, black)
    add_rect(slide, inner_x, y+194, inner_w, 108, fill=sky2, outline=line, radius=True)
    add_textbox(slide, inner_x+20, y+208, 220, 22, label(lang, "advantages"), 10.5, True, secondary)
    add_textbox(slide, inner_x+22, y+236, inner_w-44, 56, "\n".join(["▸ "+str(v) for v in data.get("technical_advantages",[])[:2]]), 8.2, False, primary)

    # IP full width
    y=1252; ip=normalize_ip(data.get("ip",{})); ip_h=220
    add_rect(slide, X, y, CW, ip_h, fill=sky2, outline=line, radius=True)
    add_textbox(slide, X+36, y+30, 260, 34, label(lang, "ip"), 15, True, primary)
    table_x, table_y, table_w, table_h = X+28, y+72, CW-56, 132
    add_rect(slide, table_x, table_y, table_w, table_h, fill=(255,255,255), outline=line)
    add_rect(slide, table_x, table_y, table_w, 40, fill=table_header, outline=line)
    col1=int(table_w*0.45); col2=int(table_w*0.28); col3=table_w-col1-col2
    for xx in [table_x+col1, table_x+col1+col2]:
        shp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(xx), px(table_y), px(1), px(table_h))
        shp.fill.solid(); shp.fill.fore_color.rgb=RGBColor(*line); shp.line.fill.background()
    add_textbox(slide, table_x+8, table_y+7, col1-16, 26, label(lang, "invention"), 11, True, black, align=PP_ALIGN.CENTER)
    add_textbox(slide, table_x+col1+8, table_y+3, col2-16, 34, ip_table_label(lang, "app_no", ip), 8.3, True, black, align=PP_ALIGN.CENTER)
    add_textbox(slide, table_x+col1+col2+8, table_y+3, col3-16, 34, ip_table_label(lang, "app_date", ip), 8.3, True, black, align=PP_ALIGN.CENTER)
    add_textbox(slide, table_x+12, table_y+48, col1-24, 56, ip["title"] or data.get("original_title",""), 8.4, False, black, align=PP_ALIGN.CENTER)
    add_textbox(slide, table_x+col1+8, table_y+48, col2-16, 56, f"{ip['application_number']}\n({ip['registration_number']})" if ip['registration_number'] else ip['application_number'], 10, False, black, align=PP_ALIGN.CENTER)
    add_textbox(slide, table_x+col1+col2+8, table_y+48, col3-16, 56, f"{ip['application_date']}\n({ip['registration_date']})" if ip['registration_date'] else ip['application_date'], 10, False, black, align=PP_ALIGN.CENTER)

    # Contact
    y=1504
    add_rect(slide, X, y, CW, 78, fill=(255,255,255), outline=line, radius=True)
    add_textbox(slide, X+42, y+24, 110, 30, label(lang, "contact"), 15.5, True, primary)
    add_textbox(slide, X+160, y+28, CW-190, 28, contact, 10, False, black)

    bio=BytesIO(); prs.save(bio); return bio.getvalue()



def _ensure_len(seq, n, default):
    out = list(seq or [])[:n]
    while len(out) < n:
        out.append(default() if callable(default) else default)
    return out

def build_data_from_edit_form(base: Dict[str, Any], university: str, department: str, professor: str, vals: Dict[str, Any]) -> Dict[str, Any]:
    """Streamlit 입력 폼 값을 내부 JSON 구조로 재조립한다."""
    return {
        "marketing_title": vals.get("marketing_title", ""),
        "subtitle": vals.get("subtitle", ""),
        "original_title": vals.get("original_title", ""),
        "university": university,
        "department": department,
        "professor": professor,
        "language": base.get("language", vals.get("language", "ko")),
        "applications": [
            {"name": vals.get(f"app_name_{i}", ""), "description": vals.get(f"app_desc_{i}", "")}
            for i in range(3)
        ],
        "overview": [vals.get(f"overview_{i}", "") for i in range(3) if vals.get(f"overview_{i}", "").strip()],
        "differentiation": [vals.get(f"diff_{i}", "") for i in range(3) if vals.get(f"diff_{i}", "").strip()],
        "limitations": [vals.get(f"limit_{i}", "") for i in range(2) if vals.get(f"limit_{i}", "").strip()],
        "technical_advantages": [vals.get(f"adv_{i}", "") for i in range(2) if vals.get(f"adv_{i}", "").strip()],
        "ip": {
            "title": vals.get("ip_title", ""),
            "application_number": vals.get("ip_application_number", ""),
            "registration_number": vals.get("ip_registration_number", ""),
            "application_date": vals.get("ip_application_date", ""),
            "registration_date": vals.get("ip_registration_date", ""),
            "applicant": vals.get("ip_applicant", ""),
        },
        "market_info": base.get("market_info", {}),
    }

# -----------------------------------------------------
# Streamlit UI
# -----------------------------------------------------
st.title("PIUM SMK 생성기")
st.caption("특허 명세서 PDF를 업로드하면 카드형 1페이지 SMK를 생성합니다.")

with st.sidebar:
    st.header("입력 정보")
    uploaded_pdf = st.file_uploader("특허 명세서 PDF 업로드", type=["pdf"])
    selected_univ = st.selectbox("대학교", UNIVERSITIES, index=0)
    custom_univ = ""
    if selected_univ == "수기입력":
        custom_univ = st.text_input("대학교 수기입력", placeholder="예: ○○대학교")
    university = custom_univ.strip() if selected_univ == "수기입력" else selected_univ
    department = st.text_input("학과/소속", placeholder="예: 활빈당공학과")
    professor = st.text_input("교수명", placeholder="예: 홍길동")
    st.caption("출원번호·출원일자는 PDF에서 확인되지 않을 때만 아래 수기입력값을 사용합니다.")
    manual_application_number = st.text_input("출원번호 (수기입력)", placeholder="예: 10-2026-0000000")
    manual_application_date = st.date_input("출원일자 (수기입력)", value=None, format="YYYY-MM-DD")
    output_language_label = st.selectbox("SMK 언어", list(LANGUAGE_OPTIONS.keys()), index=0, help="교수명/담당자 성명은 입력한 그대로 유지됩니다.")
    output_language = get_lang_code(output_language_label)

    st.divider()
    st.subheader("문의처")
    org = st.text_input("소속", placeholder="예: 부산대학교 산학협력단")
    name = st.text_input("이름", placeholder="예: 고길동")
    position = st.text_input("직책", placeholder="예: 부장")
    phone = st.text_input("연락처", placeholder="예: 051.510.2741")
    email = st.text_input("이메일", placeholder="예: example@pusan.ac.kr")

    st.divider()
    make_images = st.checkbox("적용분야 이미지 생성", value=True, help="끄면 이미지 생성 비용이 발생하지 않습니다.")
    use_logos = st.checkbox("상단 로고 자동 삽입", value=True, help="logo.zip 안의 대학 로고와 PIUM 로고를 사용합니다.")
    generate_btn = st.button("SMK 생성", type="primary", use_container_width=True)

for key, default in {
    "data": None, "brief_image": None, "pdf_bytes": None, "pptx_bytes": None,
    "pdf_path": None, "app_imgs": [], "rep_img": None, "qr_img": None, "edit_version": 0,
    "hq_image": None, "hq_png_bytes": None, "hq_pdf_bytes": None
}.items():
    if key not in st.session_state:
        st.session_state[key] = default

if generate_btn:
    if uploaded_pdf is None:
        st.error("특허 명세서 PDF를 업로드하세요."); st.stop()
    if not university:
        st.error("대학교명을 입력 또는 선택하세요."); st.stop()

    with st.spinner("PDF 분석 중..."):
        pdf_path = save_uploaded_file(uploaded_pdf)
        page_texts, ocr_used = extract_patent_pages_text(pdf_path)
        patent_text = "\n".join(page_texts)[:60000]
        authoritative_meta = extract_authoritative_application_metadata(page_texts)

        manual_no = (manual_application_number or "").strip()
        if manual_no and not APPLICATION_NO_RE.fullmatch(manual_no):
            st.error("수기입력 출원번호는 00-0000-0000000 형식으로 입력해주세요.")
            st.stop()
        manual_date = manual_application_date.strftime("%Y.%m.%d") if manual_application_date else ""

        final_app_no = authoritative_meta.get("application_number", "") or manual_no
        final_app_date = authoritative_meta.get("application_date", "") or manual_date
        missing = []
        if not final_app_no:
            missing.append("출원번호")
        if not final_app_date:
            missing.append("출원일자")
        if missing:
            st.error("출원번호 및 출원일자를 확인할 수 없습니다. 왼쪽 입력정보에서 직접 입력해주세요.")
            st.stop()

        rep_img = extract_representative_drawing(pdf_path)
        qr_img = extract_qr_code_from_first_page(pdf_path)
        st.session_state.pdf_path = pdf_path
        st.session_state.rep_img = rep_img
        st.session_state.qr_img = qr_img

    with st.spinner("GPT로 SMK 텍스트 생성 중..."):
        data = analyze_patent_with_gpt(patent_text, university, department, professor, output_language)
        data["university"] = university; data["department"] = department; data["professor"] = professor; data["language"] = output_language
        # 출원번호/출원일자는 GPT 추정값을 사용하지 않고 코드에서 확정된 권위값으로 강제한다.
        data.setdefault("ip", {})
        data["ip"]["application_number"] = final_app_no
        data["ip"]["application_date"] = final_app_date
        # 등록공보가 아닌 공개공보/출원서에서는 등록정보를 추정값으로 허용하지 않는다.
        if not authoritative_meta.get("is_registration_gazette"):
            data["ip"]["registration_number"] = ""
            data["ip"]["registration_date"] = ""
        data = localize_smk_data(data, output_language, university, department, professor)
        if ocr_used:
            st.info("텍스트 레이어가 없는 스캔 PDF로 확인되어 OCR 결과를 사용했습니다.")

    with st.spinner("시장현황 검색 및 그래프 구성 중..."):
        data["market_info"] = generate_market_info_with_web(data)
        st.session_state.data = data
        st.session_state.edit_version += 1
    university_logo = get_logo_image(university) if use_logos else None

    app_imgs=[]
    if make_images:
        with st.spinner("적용분야 이미지 세트 생성 중..."):
            try:
                app_imgs = generate_application_images_set(data.get("applications", [])[:3], university_logo=university_logo)
            except Exception as e:
                st.warning(f"적용분야 이미지 세트 생성 실패: {e}")
                # 세트 생성 실패 시 개별 생성으로 fallback
                for app in data.get("applications", [])[:3]:
                    try:
                        app_imgs.append(generate_application_image(app.get("name",""), app.get("description",""), university_logo=university_logo))
                    except Exception as e2:
                        st.warning(f"적용분야 이미지 생성 실패: {e2}")
                        app_imgs.append(Image.new("RGB", (1024,1024), "white"))
    else:
        app_imgs=[Image.new("RGB", (1024,1024), "white") for _ in data.get("applications", [])[:3]]
    st.session_state.app_imgs = app_imgs

    with st.spinner("PDF/PPTX 구성 중..."):
        contact = build_contact_text(org, name, position, phone, email, output_language)
        pium_logo = get_pium_logo_image() if use_logos else None
        piumlink_logo = get_piumlink_logo_image() if use_logos else None
        brief = compose_tech_brief(data, rep_img, app_imgs, contact, university_logo, pium_logo, st.session_state.qr_img, piumlink_logo)
        st.session_state.brief_image = brief
        st.session_state.pdf_bytes = make_pdf_bytes_from_image(brief)
        st.session_state.pptx_bytes = make_pptx_bytes(data, rep_img, app_imgs, contact, university_logo, pium_logo, st.session_state.qr_img, piumlink_logo)
        st.session_state.hq_image = None
        st.session_state.hq_png_bytes = None
        st.session_state.hq_pdf_bytes = None

if st.session_state.data is None:
    st.info("왼쪽에서 정보를 입력하고 특허 PDF를 업로드한 뒤 'SMK 생성'을 누르세요.")
else:
    col1, col2 = st.columns([1.15, 0.85], gap="large")
    with col1:
        st.subheader("SMK 미리보기")
        st.image(st.session_state.brief_image, use_container_width=True)

        if st.button("프리미엄 이미지/PDF로 변환", use_container_width=True):
            with st.spinner("레퍼런스 기반 프리미엄 인포그래픽을 생성 중..."):
                university_logo = get_logo_image(university) if use_logos else None
                pium_logo = get_pium_logo_image() if use_logos else None
                piumlink_logo = get_piumlink_logo_image() if use_logos else None
                hq_image = generate_reference_based_premium_infographic(
                    st.session_state.data,
                    st.session_state.rep_img,
                    st.session_state.app_imgs,
                    build_contact_text(org, name, position, phone, email, output_language),
                    university_logo=university_logo,
                    pium_logo=pium_logo,
                    qr_img=st.session_state.qr_img,
                    piumlink_logo=piumlink_logo,
                    base_smk_img=st.session_state.brief_image,
                )
                st.session_state.hq_image = hq_image
                st.session_state.hq_png_bytes = make_png_bytes_from_image(hq_image)
                st.session_state.hq_pdf_bytes = make_pdf_bytes_from_hq_image(hq_image)

        if st.session_state.hq_image is not None:
            st.markdown("#### 프리미엄 인포그래픽 이미지")
            st.caption("번들된 premium_infographic_reference.png를 gpt-image-2로 직접 참고해 프리미엄 인포그래픽을 생성한 결과입니다. 프리미엄 레퍼런스, 현재 SMK, 대학 로고, 대학 로고 기반 색상 팔레트, PIUM+QR 카드, 대표도면을 참조 이미지로 전달해 전체를 프리미엄 스타일로 다시 구성합니다. 별도의 후합성 없이 생성 결과 자체를 그대로 사용하며, 로고와 대표도면은 참조 이미지 기반으로 최대한 안정적으로 보존하도록 유도합니다. 생성 실패 시 기존 규칙 기반 고품질 렌더러로 자동 fallback됩니다.")
            st.image(st.session_state.hq_image, use_container_width=True)
            d1, d2 = st.columns(2)
            with d1:
                st.download_button(
                    "프리미엄 PNG 다운로드",
                    st.session_state.hq_png_bytes,
                    build_export_basename(st.session_state.data)+"_infographic.png",
                    "image/png",
                    use_container_width=True,
                )
            with d2:
                st.download_button(
                    "프리미엄 PDF 다운로드",
                    st.session_state.hq_pdf_bytes,
                    build_export_basename(st.session_state.data)+"_infographic.pdf",
                    "application/pdf",
                    use_container_width=True,
                )

        c1, c2 = st.columns(2)
        with c1:
            st.download_button("PDF 다운로드", st.session_state.pdf_bytes, build_export_basename(st.session_state.data)+".pdf", "application/pdf", type="primary", use_container_width=True)
        with c2:
            st.download_button("PPTX 다운로드(수정용)", st.session_state.pptx_bytes, build_export_basename(st.session_state.data)+".pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)

    with col2:
        st.subheader("생성 텍스트 수정")
        st.caption("JSON을 직접 고치지 않고 항목별 입력칸에서 수정할 수 있습니다.")

        data_now = st.session_state.data or {}
        apps_now = _ensure_len(data_now.get("applications", []), 3, lambda: {"name": "", "description": ""})
        overview_now = _ensure_len(data_now.get("overview", []), 3, "")
        diff_now = _ensure_len(data_now.get("differentiation", []), 3, "")
        limit_now = _ensure_len(data_now.get("limitations", []), 2, "")
        adv_now = _ensure_len(data_now.get("technical_advantages", []), 2, "")
        ip_now = normalize_ip(data_now.get("ip", {}))
        k = f"edit_v{st.session_state.edit_version}"

        with st.form("edit_text_form"):
            st.markdown("#### 기본 정보")
            marketing_title = st.text_input("기술명", value=data_now.get("marketing_title", ""), key=f"{k}_marketing_title")
            subtitle = st.text_input("한 줄 요약", value=data_now.get("subtitle", ""), key=f"{k}_subtitle")
            original_title = st.text_input("원 발명의 명칭", value=data_now.get("original_title", ""), key=f"{k}_original_title")

            st.markdown("#### 적용분야 / 제품")
            app_vals = []
            for i in range(3):
                with st.expander(f"적용분야 {i+1}", expanded=True):
                    app_name = st.text_input(f"적용분야 {i+1} 명칭", value=apps_now[i].get("name", ""), key=f"{k}_app_name_{i}")
                    app_desc = st.text_area(f"적용분야 {i+1} 설명", value=apps_now[i].get("description", ""), height=70, key=f"{k}_app_desc_{i}")
                    app_vals.append((app_name, app_desc))

            st.markdown("#### 기술개요")
            overview_vals = []
            for i in range(3):
                overview_vals.append(st.text_area(f"기술개요 {i+1}", value=overview_now[i], height=75, key=f"{k}_overview_{i}"))

            st.markdown("#### 핵심 차별성")
            diff_vals = []
            for i in range(3):
                diff_vals.append(st.text_area(f"핵심 차별성 {i+1}", value=diff_now[i], height=75, key=f"{k}_diff_{i}"))

            st.markdown("#### 기술 경쟁력")
            st.caption("기존기술 한계")
            limit_vals = []
            for i in range(2):
                limit_vals.append(st.text_area(f"기존기술 한계 {i+1}", value=limit_now[i], height=70, key=f"{k}_limit_{i}"))
            st.caption("기술적 우위")
            adv_vals = []
            for i in range(2):
                adv_vals.append(st.text_area(f"기술적 우위 {i+1}", value=adv_now[i], height=70, key=f"{k}_adv_{i}"))

            st.caption("※ 적용분야를 수정하면 시장현황도 함께 재검색됩니다.")

            st.markdown("#### 지식재산권 현황")
            ip_title = st.text_input("발명의 명칭", value=ip_now.get("title", ""), key=f"{k}_ip_title")
            c_ip1, c_ip2 = st.columns(2)
            with c_ip1:
                ip_application_number = st.text_input("출원번호", value=ip_now.get("application_number", ""), key=f"{k}_ip_application_number")
                ip_application_date = st.text_input("출원일자", value=ip_now.get("application_date", ""), key=f"{k}_ip_application_date")
            with c_ip2:
                ip_registration_number = st.text_input("등록번호", value=ip_now.get("registration_number", ""), key=f"{k}_ip_registration_number")
                ip_registration_date = st.text_input("등록일자", value=ip_now.get("registration_date", ""), key=f"{k}_ip_registration_date")
            ip_applicant = st.text_input("출원인", value=ip_now.get("applicant", ""), key=f"{k}_ip_applicant")

            submitted = st.form_submit_button("수정 내용으로 PDF/PPTX 다시 생성", use_container_width=True)

        if submitted:
            vals = {
                "marketing_title": marketing_title,
                "subtitle": subtitle,
                "original_title": original_title,
                "ip_title": ip_title,
                "ip_application_number": ip_application_number,
                "ip_registration_number": ip_registration_number,
                "ip_application_date": ip_application_date,
                "ip_registration_date": ip_registration_date,
                "ip_applicant": ip_applicant,
            }
            for i, (app_name, app_desc) in enumerate(app_vals):
                vals[f"app_name_{i}"] = app_name
                vals[f"app_desc_{i}"] = app_desc
            for i, v in enumerate(overview_vals):
                vals[f"overview_{i}"] = v
            for i, v in enumerate(diff_vals):
                vals[f"diff_{i}"] = v
            for i, v in enumerate(limit_vals):
                vals[f"limit_{i}"] = v
            for i, v in enumerate(adv_vals):
                vals[f"adv_{i}"] = v

            with st.spinner("수정 내용과 시장현황을 반영하는 중..."):
                edited = build_data_from_edit_form(data_now, university, department, professor, vals)
                edited["language"] = output_language
                edited = localize_smk_data(edited, output_language, university, department, professor)
                edited["market_info"] = generate_market_info_with_web(edited)
                contact = build_contact_text(org, name, position, phone, email, output_language)
                rep_img = st.session_state.rep_img or extract_representative_drawing(st.session_state.pdf_path)
                university_logo = get_logo_image(university) if use_logos else None
                pium_logo = get_pium_logo_image() if use_logos else None
                piumlink_logo = get_piumlink_logo_image() if use_logos else None
                brief = compose_tech_brief(edited, rep_img, st.session_state.app_imgs, contact, university_logo, pium_logo, st.session_state.qr_img, piumlink_logo)
                st.session_state.data = edited
                st.session_state.brief_image = brief
                st.session_state.pdf_bytes = make_pdf_bytes_from_image(brief)
                st.session_state.pptx_bytes = make_pptx_bytes(edited, rep_img, st.session_state.app_imgs, contact, university_logo, pium_logo, st.session_state.qr_img, piumlink_logo)
                st.session_state.hq_image = None
                st.session_state.hq_png_bytes = None
                st.session_state.hq_pdf_bytes = None
                st.session_state.edit_version += 1
            st.success("수정 내용이 반영되었습니다.")
            st.rerun()

        with st.expander("고급 사용자용 JSON 보기"):
            st.json(st.session_state.data)